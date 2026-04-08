"""Extract embedded Picture images from .pptx files."""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass(frozen=True)
class ExtractedImage:
    slide_index: int  # 1-based for human output
    image_index: int  # 1-based per slide
    blob: bytes
    mime_type: str


def _iter_pictures(shape, slide_idx: int, counter: list[int]):
    """Recursively walk shapes, yielding pictures (handles groups)."""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        counter[0] += 1
        img = shape.image
        yield ExtractedImage(
            slide_index=slide_idx,
            image_index=counter[0],
            blob=img.blob,
            mime_type=img.content_type,
        )
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from _iter_pictures(child, slide_idx, counter)


def _parse_slide_range(slide_range: str, total: int) -> set[int]:
    """Parse '1-3,5' style range; 'all' returns every 1-based slide index."""
    if slide_range == "all":
        return set(range(1, total + 1))
    result: set[int] = set()
    for chunk in slide_range.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        if "-" in chunk:
            lo, hi = chunk.split("-", 1)
            result.update(range(int(lo), int(hi) + 1))
        else:
            result.add(int(chunk))
    return {i for i in result if 1 <= i <= total}


def extract_images(
    pptx_path: Path,
    slide_range: str = "all",
) -> list[ExtractedImage]:
    """Return all Picture shapes in the selected slides, in reading order."""
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    wanted = _parse_slide_range(slide_range, total)

    out: list[ExtractedImage] = []
    for idx, slide in enumerate(prs.slides, start=1):
        if idx not in wanted:
            continue
        counter = [0]
        for shape in slide.shapes:
            out.extend(_iter_pictures(shape, idx, counter))
    return out
