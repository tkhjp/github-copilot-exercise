#!/usr/bin/env python
"""Generate PPTX mirrors of the judgment test cases (tc02, tc03).

These files are the `copilot_pptx` arm of the Phase 4 text-vs-image extension:
- tc02 — UI Change Request mockup + 4 red callouts
- tc03 — AWS multi-AZ architecture with VPC boundary, subnets, colored links

The pixel layout mirrors `generate_test_images.py` 1:1 via a 96-DPI pixel→EMU
mapping, so the PPTX visually corresponds to the PNG while giving Copilot the
option to read native shape text instead of OCR.

Run:
    python tests/text_vs_image/generate_test_pptx.py
    # writes tests/text_vs_image/inputs/{02_ui_change,03_complex_arch}.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parent
INPUTS = ROOT / "inputs"


def px(n: float) -> Emu:
    """Convert a 96-DPI pixel coordinate to EMU so PIL layouts port over 1:1."""
    return Emu(int(n * 9525))


def _hex(s: str) -> RGBColor:
    return RGBColor.from_string(s.lstrip("#"))


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _add_box(
    slide,
    x, y, w, h,
    fill=None,
    outline=None,
    outline_w=1.0,
    shape=MSO_SHAPE.RECTANGLE,
):
    shp = slide.shapes.add_shape(shape, px(x), px(y), px(w), px(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex(fill)
    if outline is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _hex(outline)
        shp.line.width = Pt(outline_w)
    tf = shp.text_frame
    tf.text = ""
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return shp


def _add_text(
    slide,
    x, y, w, h,
    text,
    size_pt=11,
    bold=False,
    color="#000000",
    align=PP_ALIGN.LEFT,
    font_name="Arial",
):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = _hex(color)
    return tb


def _text_in_shape(shp, text, size_pt=11, bold=False, color="#000000", align=PP_ALIGN.CENTER):
    tf = shp.text_frame
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = _hex(color)


def _add_line(slide, x1, y1, x2, y2, color="#000000", width_pt=2.0, dashed=False):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2))
    ln.line.color.rgb = _hex(color)
    ln.line.width = Pt(width_pt)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return ln


def _add_dashed_rect(slide, x, y, w, h, color="#7c3aed", width_pt=2.0):
    for (x1, y1, x2, y2) in [
        (x, y, x + w, y),
        (x, y + h, x + w, y + h),
        (x, y, x, y + h),
        (x + w, y, x + w, y + h),
    ]:
        _add_line(slide, x1, y1, x2, y2, color=color, width_pt=width_pt, dashed=True)


# ---------------------------------------------------------------------------
# tc02 — UI Change Request
# ---------------------------------------------------------------------------

def build_tc02(path: Path) -> None:
    W, H = 1280, 820
    prs = Presentation()
    prs.slide_width = px(W)
    prs.slide_height = px(H)
    slide = _blank_slide(prs)

    # Document header
    _add_text(slide, 24, 12, 900, 28,
              "UI Change Request — ShopApp Product Search",
              size_pt=18, bold=True)
    _add_text(slide, 24, 44, 900, 22,
              "Reviewer: PM / Design — 2026-04-10",
              size_pt=12, color="#6b7280")
    _add_line(slide, 24, 78, W - 24, 78, color="#e5e7eb", width_pt=1.5)

    # Screen mockup frame
    sx1, sy1, sx2, sy2 = 40, 90, 780, 760
    _add_box(slide, sx1, sy1, sx2 - sx1, sy2 - sy1,
             fill="#ffffff", outline="#9ca3af", outline_w=1.5)

    # Dark header bar inside screen
    hdr_h = 50
    _add_box(slide, sx1, sy1, sx2 - sx1, hdr_h, fill="#1e293b", outline="#1e293b")
    _add_text(slide, sx1 + 20, sy1 + 13, 200, 28,
              "ShopApp", size_pt=16, bold=True, color="#ffffff")
    # Login button
    login_x1 = sx2 - 100
    _add_box(slide, login_x1, sy1 + 12, 80, 26,
             fill="#374151", outline="#9ca3af", outline_w=0.75)
    _add_text(slide, login_x1, sy1 + 14, 80, 22,
              "Login", size_pt=11, color="#ffffff", align=PP_ALIGN.CENTER)

    # Search bar
    search_y = sy1 + hdr_h + 20
    sb_x1 = sx1 + 20
    sb_x2 = sx2 - 140
    _add_box(slide, sb_x1, search_y, sb_x2 - sb_x1, 36,
             fill="#ffffff", outline="#9ca3af", outline_w=1.5)
    _add_text(slide, sb_x1 + 10, search_y + 8, 300, 22,
              "Search products...", size_pt=11, color="#9ca3af")
    # Blue search button
    btn_x1 = sb_x2 + 10
    _add_box(slide, btn_x1, search_y, 100, 36,
             fill="#2563eb", outline="#1e40af", outline_w=1.5)
    _add_text(slide, btn_x1, search_y + 8, 100, 22,
              "Search", size_pt=12, bold=True, color="#ffffff", align=PP_ALIGN.CENTER)

    # Filter sidebar
    fx1 = sx1 + 20
    fx2 = fx1 + 180
    fy1 = search_y + 70
    fy2 = fy1 + 290
    _add_box(slide, fx1, fy1, fx2 - fx1, fy2 - fy1,
             fill="#f9fafb", outline="#d1d5db", outline_w=0.75)
    _add_text(slide, fx1 + 12, fy1 + 10, 160, 22,
              "Filters", size_pt=13, bold=True)
    for i, cat in enumerate(["Electronics", "Books", "Clothing"]):
        cy = fy1 + 50 + i * 32
        _add_box(slide, fx1 + 14, cy, 16, 16,
                 fill="#ffffff", outline="#6b7280", outline_w=1.25)
        _add_text(slide, fx1 + 38, cy - 2, 120, 22,
                  cat, size_pt=11)
    _add_text(slide, fx1 + 12, fy1 + 160, 160, 22,
              "Price range", size_pt=11)
    # slider track
    sl_y = fy1 + 200
    _add_line(slide, fx1 + 14, sl_y, fx2 - 14, sl_y, color="#9ca3af", width_pt=2.5)
    # two knobs
    _add_box(slide, fx1 + 30, sl_y - 6, 12, 12,
             fill="#2563eb", outline="#1e40af", shape=MSO_SHAPE.OVAL, outline_w=1.0)
    _add_box(slide, fx2 - 42, sl_y - 6, 12, 12,
             fill="#2563eb", outline="#1e40af", shape=MSO_SHAPE.OVAL, outline_w=1.0)
    _add_text(slide, fx1 + 12, fy1 + 215, 160, 18,
              "¥0 - ¥50,000", size_pt=10, color="#6b7280")

    # Results table
    tx1 = fx2 + 20
    tx2 = sx2 - 20
    ty1 = search_y + 70
    row_h = 40
    # header
    _add_box(slide, tx1, ty1, tx2 - tx1, row_h,
             fill="#e5e7eb", outline="#9ca3af", outline_w=0.75)
    cols = [("Product", 220), ("Price", 110), ("Stock", 90)]
    cx = tx1
    for name, cw in cols:
        _add_text(slide, cx + 10, ty1 + 10, cw - 10, 22,
                  name, size_pt=11, bold=True)
        cx += cw
    products = [
        ("Wireless Headphones", "¥12,800", "42"),
        ("USB-C Hub (7-in-1)", "¥4,500", "118"),
        ("Mech. Keyboard", "¥15,200", "7"),
    ]
    for r, row in enumerate(products):
        ry = ty1 + row_h * (r + 1)
        bg = "#ffffff" if r % 2 == 0 else "#f9fafb"
        _add_box(slide, tx1, ry, tx2 - tx1, row_h,
                 fill=bg, outline="#e5e7eb", outline_w=0.75)
        cx = tx1
        for val, (_n, cw) in zip(row, cols):
            _add_text(slide, cx + 10, ry + 10, cw - 10, 22,
                      val, size_pt=11)
            cx += cw

    # Pagination
    page_y = ty1 + row_h * (len(products) + 1) + 20
    page_cx = (tx1 + tx2) // 2
    labels = ["<", "1", "2", "3", ">"]
    btn_w = 34
    btn_h = 30
    total_w = len(labels) * (btn_w + 6) - 6
    start_x = page_cx - total_w // 2
    for i, lbl in enumerate(labels):
        bx = start_x + i * (btn_w + 6)
        is_curr = (lbl == "1")
        _add_box(slide, bx, page_y, btn_w, btn_h,
                 fill="#2563eb" if is_curr else "#ffffff",
                 outline="#9ca3af", outline_w=0.75)
        _add_text(slide, bx, page_y + 6, btn_w, 22,
                  lbl, size_pt=11,
                  color="#ffffff" if is_curr else "#1e293b",
                  align=PP_ALIGN.CENTER)

    # ------ 4 red callouts ------
    callouts = [
        (
            (sb_x1 + 150, search_y + 40),
            '1. Add a "Remember last search" checkbox below the search bar.',
            820, 120,
        ),
        (
            (btn_x1 + 50, search_y + 18),
            "2. Change the Search button color from blue to green.",
            820, 230,
        ),
        (
            (tx2 - 20, ty1 + 15),
            "3. Add a CSV export button on the results table top-right corner.",
            820, 360,
        ),
        (
            (page_cx, page_y + 15),
            "4. Show 20 items per page instead of the current 10.",
            820, 490,
        ),
    ]
    box_w = 260
    box_h = 95
    for (target, text, bx, by) in callouts:
        shp = _add_box(slide, bx, by, box_w, box_h,
                       fill="#fef2f2", outline="#dc2626", outline_w=1.75,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text_in_shape(shp, text, size_pt=11, bold=True, color="#991b1b", align=PP_ALIGN.LEFT)
        shp.text_frame.margin_left = Emu(px(10))
        shp.text_frame.margin_right = Emu(px(10))
        shp.text_frame.margin_top = Emu(px(8))
        shp.text_frame.margin_bottom = Emu(px(8))
        # red leader line box-center -> target
        cx = bx + box_w / 2
        cy = by + box_h / 2
        _add_line(slide, cx, cy, target[0], target[1], color="#dc2626", width_pt=1.75)

    # Legend
    lg_y = 650
    _add_box(slide, 820, lg_y, 240, 90,
             fill="#fefce8", outline="#ca8a04", outline_w=0.75)
    _add_text(slide, 832, lg_y + 8, 200, 22,
              "Legend", size_pt=12, bold=True, color="#854d0e")
    _add_line(slide, 832, lg_y + 40, 862, lg_y + 40, color="#dc2626", width_pt=1.75)
    _add_text(slide, 870, lg_y + 30, 180, 20,
              "Red = change request", size_pt=10, color="#854d0e")
    _add_box(slide, 832, lg_y + 56, 30, 18,
             fill="#ffffff", outline="#9ca3af", outline_w=0.75)
    _add_text(slide, 870, lg_y + 58, 180, 20,
              "Gray = existing UI", size_pt=10, color="#854d0e")

    prs.save(str(path))


# ---------------------------------------------------------------------------
# tc03 — AWS Multi-AZ Architecture
# ---------------------------------------------------------------------------

def _add_rounded_component(slide, rect, label, fill, outline="#000000"):
    x1, y1, x2, y2 = rect
    shp = _add_box(slide, x1, y1, x2 - x1, y2 - y1,
                   fill=fill, outline=outline, outline_w=1.5,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _text_in_shape(shp, label, size_pt=11, bold=True, color="#000000")
    return shp


def _add_rect_component(slide, rect, label, fill, outline="#000000"):
    x1, y1, x2, y2 = rect
    shp = _add_box(slide, x1, y1, x2 - x1, y2 - y1,
                   fill=fill, outline=outline, outline_w=1.5)
    _text_in_shape(shp, label, size_pt=11, bold=True, color="#000000")
    return shp


def _add_cylinder_component(slide, rect, label, fill, outline="#000000"):
    """Approximate the PIL cylinder with a rounded rectangle — pptx has MSO_SHAPE.CAN
    but tends to be awkwardly tall; rounded rectangle with label is close enough
    visually and carries the same semantics (it's a database)."""
    x1, y1, x2, y2 = rect
    shp = slide.shapes.add_shape(MSO_SHAPE.CAN, px(x1), px(y1), px(x2 - x1), px(y2 - y1))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex(fill)
    shp.line.color.rgb = _hex(outline)
    shp.line.width = Pt(1.5)
    _text_in_shape(shp, label, size_pt=11, bold=True, color="#000000")
    return shp


def build_tc03(path: Path) -> None:
    W, H = 1400, 900
    prs = Presentation()
    prs.slide_width = px(W)
    prs.slide_height = px(H)
    slide = _blank_slide(prs)

    # Title
    _add_text(slide, 30, 16, 900, 32,
              "Production Architecture — AWS Multi-AZ",
              size_pt=20, bold=True)
    _add_text(slide, 30, 50, 900, 22,
              "Environment: prod • Region: ap-northeast-1",
              size_pt=11, color="#6b7280")

    # --- Top row: Route53, CloudFront, S3 on left; IAM, CloudWatch on right ---
    top_y = 100
    top_h = 60
    route53 = (50, top_y, 180, top_y + top_h)
    cloudfront = (230, top_y, 380, top_y + top_h)
    s3 = (430, top_y, 560, top_y + top_h)
    iam = (1230, top_y, 1360, top_y + top_h)
    cw_rect = (1230, top_y + 80, 1360, top_y + 80 + top_h)

    _add_rounded_component(slide, route53, "Route 53\nDNS", "#e9d5ff")
    _add_rounded_component(slide, cloudfront, "CloudFront\nCDN", "#fde68a")
    _add_rect_component(slide, s3, "S3\nStatic Assets", "#d1fae5")
    _add_rounded_component(slide, iam, "IAM\nRoles", "#fbcfe8")
    _add_rounded_component(slide, cw_rect, "CloudWatch\nLogs/Metrics", "#fbcfe8")

    # --- VPC boundary (purple dashed) ---
    vpc = (50, 210, 1200, 820)
    _add_dashed_rect(slide, vpc[0], vpc[1], vpc[2] - vpc[0], vpc[3] - vpc[1],
                     color="#7c3aed", width_pt=2.25)
    _add_text(slide, vpc[0] + 12, vpc[1] + 6, 260, 22,
              "VPC: 10.0.0.0/16", size_pt=12, bold=True, color="#7c3aed")

    # --- AZ boundaries (cyan dashed) ---
    az_a = (75, 260, 620, 800)
    az_b = (640, 260, 1180, 800)
    for az in (az_a, az_b):
        _add_dashed_rect(slide, az[0], az[1], az[2] - az[0], az[3] - az[1],
                         color="#0891b2", width_pt=1.5)
    _add_text(slide, az_a[0] + 10, az_a[1] + 4, 300, 22,
              "AZ-a (ap-northeast-1a)", size_pt=11, bold=True, color="#0891b2")
    _add_text(slide, az_b[0] + 10, az_b[1] + 4, 300, 22,
              "AZ-b (ap-northeast-1c)", size_pt=11, bold=True, color="#0891b2")

    # --- Public Subnets (blue background) ---
    pub_a = (95, 295, 600, 395)
    pub_b = (660, 295, 1160, 395)
    for r in (pub_a, pub_b):
        _add_box(slide, r[0], r[1], r[2] - r[0], r[3] - r[1],
                 fill="#dbeafe", outline="#3b82f6", outline_w=0.75)
    _add_text(slide, pub_a[0] + 8, pub_a[1] + 2, 200, 18,
              "Public Subnet", size_pt=10, color="#1e3a8a")
    _add_text(slide, pub_b[0] + 8, pub_b[1] + 2, 200, 18,
              "Public Subnet", size_pt=10, color="#1e3a8a")

    alb_a = (pub_a[0] + 40, pub_a[1] + 30, pub_a[0] + 200, pub_a[1] + 85)
    alb_b = (pub_b[0] + 40, pub_b[1] + 30, pub_b[0] + 200, pub_b[1] + 85)
    _add_rounded_component(slide, alb_a, "ALB\n(AZ-a)", "#fef3c7")
    _add_rounded_component(slide, alb_b, "ALB\n(AZ-b)", "#fef3c7")

    nat_a = (pub_a[0] + 280, pub_a[1] + 30, pub_a[0] + 430, pub_a[1] + 85)
    nat_b = (pub_b[0] + 280, pub_b[1] + 30, pub_b[0] + 430, pub_b[1] + 85)
    _add_rect_component(slide, nat_a, "NAT GW\n(AZ-a)", "#fef3c7")
    _add_rect_component(slide, nat_b, "NAT GW\n(AZ-b)", "#fef3c7")

    # --- Private Subnets (yellow background) ---
    priv_a = (95, 420, 600, 570)
    priv_b = (660, 420, 1160, 570)
    for r in (priv_a, priv_b):
        _add_box(slide, r[0], r[1], r[2] - r[0], r[3] - r[1],
                 fill="#fef9c3", outline="#eab308", outline_w=0.75)
    _add_text(slide, priv_a[0] + 8, priv_a[1] + 2, 260, 18,
              "Private Subnet (App)", size_pt=10, color="#713f12")
    _add_text(slide, priv_b[0] + 8, priv_b[1] + 2, 260, 18,
              "Private Subnet (App)", size_pt=10, color="#713f12")

    ecs_a = (priv_a[0] + 40, priv_a[1] + 40, priv_a[0] + 230, priv_a[1] + 130)
    ecs_b = (priv_b[0] + 40, priv_b[1] + 40, priv_b[0] + 230, priv_b[1] + 130)
    _add_rounded_component(slide, ecs_a, "ECS Fargate\nApp Service\n(AZ-a)", "#fed7aa")
    _add_rounded_component(slide, ecs_b, "ECS Fargate\nApp Service\n(AZ-b)", "#fed7aa")

    redis_a = (priv_a[0] + 290, priv_a[1] + 50, priv_a[0] + 450, priv_a[1] + 120)
    redis_b = (priv_b[0] + 290, priv_b[1] + 50, priv_b[0] + 450, priv_b[1] + 120)
    _add_rounded_component(slide, redis_a, "ElastiCache\nRedis", "#fecaca")
    _add_rounded_component(slide, redis_b, "ElastiCache\nRedis", "#fecaca")

    # --- DB Subnets (green background) ---
    db_a = (95, 595, 600, 780)
    db_b = (660, 595, 1160, 780)
    for r in (db_a, db_b):
        _add_box(slide, r[0], r[1], r[2] - r[0], r[3] - r[1],
                 fill="#dcfce7", outline="#16a34a", outline_w=0.75)
    _add_text(slide, db_a[0] + 8, db_a[1] + 2, 200, 18,
              "DB Subnet", size_pt=10, color="#14532d")
    _add_text(slide, db_b[0] + 8, db_b[1] + 2, 200, 18,
              "DB Subnet", size_pt=10, color="#14532d")

    rds_primary = (db_a[0] + 100, db_a[1] + 50, db_a[0] + 280, db_a[1] + 150)
    rds_replica = (db_b[0] + 100, db_b[1] + 50, db_b[0] + 280, db_b[1] + 150)
    _add_cylinder_component(slide, rds_primary, "RDS PostgreSQL\nPrimary", "#bbf7d0")
    _add_cylinder_component(slide, rds_replica, "RDS PostgreSQL\nRead Replica", "#bbf7d0")

    # ------ Arrows ------
    def bottom_mid(r): return ((r[0] + r[2]) / 2, r[3])
    def top_mid(r): return ((r[0] + r[2]) / 2, r[1])
    def left_mid(r): return (r[0], (r[1] + r[3]) / 2)
    def right_mid(r): return (r[2], (r[1] + r[3]) / 2)

    # Route53 → CloudFront, CloudFront → S3 (black)
    _add_line(slide, *right_mid(route53), *left_mid(cloudfront), color="#000000", width_pt=1.75)
    _add_line(slide, *right_mid(cloudfront), *left_mid(s3), color="#000000", width_pt=1.75)
    _add_text(slide, (cloudfront[2] + s3[0]) / 2 - 20, cloudfront[1] - 16, 60, 16,
              "origin", size_pt=9, color="#6b7280")

    # CloudFront → both ALBs (black, labeled HTTPS)
    _add_line(slide, *bottom_mid(cloudfront), *top_mid(alb_a), color="#000000", width_pt=1.75)
    _add_line(slide, *bottom_mid(cloudfront), *top_mid(alb_b), color="#000000", width_pt=1.75)
    _add_text(slide, cloudfront[0] - 10, cloudfront[3] + 8, 60, 16,
              "HTTPS", size_pt=9, color="#6b7280")

    # ALB → ECS (same AZ, black)
    _add_line(slide, *bottom_mid(alb_a), *top_mid(ecs_a), color="#000000", width_pt=1.75)
    _add_line(slide, *bottom_mid(alb_b), *top_mid(ecs_b), color="#000000", width_pt=1.75)

    # ECS → ElastiCache (same AZ, black, label GET/SET)
    _add_line(slide, *right_mid(ecs_a), *left_mid(redis_a), color="#000000", width_pt=1.75)
    _add_line(slide, *right_mid(ecs_b), *left_mid(redis_b), color="#000000", width_pt=1.75)
    _add_text(slide, ecs_a[2] + 5, ecs_a[1] + 10, 80, 16,
              "GET/SET", size_pt=9, color="#6b7280")

    # ECS → RDS Primary (write) — GREEN thick line, both AZs
    _add_line(slide, *bottom_mid(ecs_a), *top_mid(rds_primary), color="#16a34a", width_pt=3.5)
    _add_line(slide, ecs_b[0] + 40, ecs_b[3], rds_primary[2] - 20, rds_primary[1],
              color="#16a34a", width_pt=2.5)
    _add_text(slide, ecs_a[2] - 30, (ecs_a[3] + rds_primary[1]) / 2 - 8, 60, 16,
              "write", size_pt=10, bold=True, color="#16a34a")

    # ECS → RDS Replica (read-only) — CYAN thick line
    _add_line(slide, *bottom_mid(ecs_b), *top_mid(rds_replica), color="#0891b2", width_pt=3.5)
    _add_text(slide, ecs_b[2] - 100, (ecs_b[3] + rds_replica[1]) / 2 - 8, 100, 16,
              "read-only", size_pt=10, bold=True, color="#0891b2")

    # NAT GW → Internet (upward, gray)
    for nat in (nat_a, nat_b):
        _add_line(slide, (nat[0] + nat[2]) / 2, nat[1],
                  (nat[0] + nat[2]) / 2, nat[1] - 40,
                  color="#6b7280", width_pt=1.75)
    _add_text(slide, nat_a[0] + 10, nat_a[1] - 58, 100, 16,
              "Outbound", size_pt=9, color="#6b7280")

    # Everything → CloudWatch (thin light gray)
    for src in [bottom_mid(ecs_a), bottom_mid(ecs_b), bottom_mid(alb_a)]:
        _add_line(slide, *src, *left_mid(cw_rect), color="#e5e7eb", width_pt=0.75)

    # Legend
    lg_x, lg_y = 50, H - 50
    _add_box(slide, lg_x, lg_y, 960, 36,
             fill="#f9fafb", outline="#d1d5db", outline_w=0.75)
    items = [
        ("#dbeafe", "Public Subnet"),
        ("#fef9c3", "Private Subnet (App)"),
        ("#dcfce7", "DB Subnet"),
        ("#e9d5ff", "External / DNS"),
        ("#fecaca", "Cache"),
    ]
    ix = lg_x + 12
    for color, label in items:
        _add_box(slide, ix, lg_y + 10, 20, 16,
                 fill=color, outline="#000000", outline_w=0.5)
        _add_text(slide, ix + 26, lg_y + 10, 180, 18,
                  label, size_pt=10)
        ix += 26 + len(label) * 7 + 24

    prs.save(str(path))


def main() -> None:
    INPUTS.mkdir(parents=True, exist_ok=True)
    tc02_path = INPUTS / "02_ui_change.pptx"
    tc03_path = INPUTS / "03_complex_arch.pptx"
    build_tc02(tc02_path)
    print(f"wrote {tc02_path} ({tc02_path.stat().st_size} bytes)")
    build_tc03(tc03_path)
    print(f"wrote {tc03_path} ({tc03_path.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
