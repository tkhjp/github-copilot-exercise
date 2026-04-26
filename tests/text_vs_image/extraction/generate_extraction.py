#!/usr/bin/env python
"""Generate PPTX + 8 PNGs + ground_truth.yaml for the extraction prompt experiment.

Single command produces all artifacts consumed by Copilot Web trials:

    python tests/text_vs_image/extraction/generate_extraction.py

All content flows from `extraction_spec.SPEC`. To change a slide, edit the spec
and re-run this script.
"""
from __future__ import annotations

import argparse
import math
import sys
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

# Allow both `python script.py` and `python -m tests.text_vs_image.extraction.generate_extraction`.
# Without this bootstrap the absolute `from tests.text_vs_image...` import below requires the repo
# root to already be on sys.path (via PYTHONPATH or pytest's pythonpath config).
_REPO_ROOT = Path(__file__).resolve().parents[3]
if str(_REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_REPO_ROOT))

from tests.text_vs_image.extraction.extraction_spec import SPEC, emit_ground_truth_yaml


ROOT = Path(__file__).resolve().parent
CANVAS_W, CANVAS_H = 1600, 900

# Canonical PNG filenames per pattern id. Consumed by main() and by the tests
# (single source of truth — don't duplicate this list).
PNG_FILENAMES = {
    "p01": "p01_ui_callouts.png",
    "p02": "p02_before_after.png",
    "p03": "p03_process_flow.png",
    "p04": "p04_dashboard_annotated.png",
    "p05": "p05_hierarchical_drilldown.png",
    "p06": "p06_review_comments.png",
    "p07": "p07_mixed_dashboard.png",
    "p08": "p08_org_chart.png",
}

# Shared visual palette (kept deliberately muted so callouts stand out).
COLORS = {
    "bg":          "#ffffff",
    "card_bg":     "#f9fafb",
    "card_border": "#d1d5db",
    "text":        "#111827",
    "muted":       "#6b7280",
    "header_bg":   "#1e293b",
    "header_text": "#ffffff",
    "primary":     "#2563eb",
    "primary_dk":  "#1e40af",
    "success":     "#16a34a",
    "warn":        "#e67e22",
    "danger":      "#dc2626",
    "danger_bg":   "#fef2f2",
    "danger_text": "#991b1b",
    "grid":        "#e5e7eb",
}


def _font(size: int) -> ImageFont.ImageFont:
    """Load a Japanese-capable font. Falls back gracefully if unavailable."""
    for path in [
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
    ]:
        try:
            return ImageFont.truetype(path, size)
        except OSError:
            continue
    return ImageFont.load_default()


def _bold_font(size: int) -> ImageFont.ImageFont:
    for path in [
        "/System/Library/Fonts/Hiragino Sans GB.ttc",  # same TTC as _font; index 1 is bold on macOS
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
    ]:
        try:
            return ImageFont.truetype(path, size, index=1 if path.endswith(".ttc") else 0)
        except OSError:
            continue
    return _font(size)


def _mono_font(size: int) -> ImageFont.ImageFont:
    for path in ["/System/Library/Fonts/Menlo.ttc", "/System/Library/Fonts/Monaco.ttf"]:
        try:
            return ImageFont.truetype(path, size)
        except OSError:
            continue
    return ImageFont.load_default()


def _tsize(draw: ImageDraw.ImageDraw, text: str, font) -> tuple[int, int]:
    bbox = draw.textbbox((0, 0), text, font=font)
    return bbox[2] - bbox[0], bbox[3] - bbox[1]


def _arrow(draw: ImageDraw.ImageDraw, p1, p2, color="black", width=2, head=12):
    draw.line([p1, p2], fill=color, width=width)
    ang = math.atan2(p2[1] - p1[1], p2[0] - p1[0])
    a1 = (p2[0] - head * math.cos(ang - math.radians(22)),
          p2[1] - head * math.sin(ang - math.radians(22)))
    a2 = (p2[0] - head * math.cos(ang + math.radians(22)),
          p2[1] - head * math.sin(ang + math.radians(22)))
    draw.polygon([p2, a1, a2], fill=color)


def _draw_screenshot_card(
    draw: ImageDraw.ImageDraw, rect: tuple[int, int, int, int], *,
    title: str | None = None, title_font=None, body_font=None,
):
    """Draw a 'screenshot card' rectangle with an optional title bar and muted border.
    Body content is drawn by the caller by reading the spec and calling further primitives.
    """
    x1, y1, x2, y2 = rect
    draw.rectangle([x1, y1, x2, y2], fill=COLORS["card_bg"], outline=COLORS["card_border"], width=2)
    if title:
        tf = title_font or _bold_font(14)
        draw.rectangle([x1, y1, x2, y1 + 32], fill=COLORS["header_bg"])
        draw.text((x1 + 12, y1 + 8), title, fill=COLORS["header_text"], font=tf)


def _draw_callout(
    draw: ImageDraw.ImageDraw, *,
    box_rect: tuple[int, int, int, int],
    label: str,
    text: str,
    target: tuple[int, int],
    font=None,
    color=None,
    fill=None,
    text_color=None,
):
    """Red callout box with label (e.g. 'C1') + text + leader line to target point."""
    x1, y1, x2, y2 = box_rect
    f = font or _bold_font(13)
    color = color or COLORS["danger"]
    fill = fill or COLORS["danger_bg"]
    text_color = text_color or COLORS["danger_text"]
    draw.rectangle([x1, y1, x2, y2], fill=fill, outline=color, width=2)
    draw.text((x1 + 10, y1 + 8), f"{label} {text}", fill=text_color, font=f)
    cx, cy = (x1 + x2) // 2, (y1 + y2) // 2
    _arrow(draw, (cx, cy), target, color=color, width=2, head=10)


def _new_canvas() -> Image.Image:
    return Image.new("RGB", (CANVAS_W, CANVAS_H), COLORS["bg"])


def _draw_slide_title(draw: ImageDraw.ImageDraw, title: str, subtitle: str | None = None) -> int:
    """Draw the slide title bar; return y-offset where body content can start."""
    draw.text((32, 22), title, fill=COLORS["text"], font=_bold_font(26))
    if subtitle:
        draw.text((32, 60), subtitle, fill=COLORS["muted"], font=_font(13))
    draw.line([(32, 86), (CANVAS_W - 32, 86)], fill=COLORS["grid"], width=2)
    return 100  # body starts here


def _draw_footer(draw: ImageDraw.ImageDraw, text: str) -> None:
    """Small muted copyright/metadata line at the very bottom of the canvas."""
    draw.text((32, CANVAS_H - 24), text, fill=COLORS["muted"], font=_font(11))


# --- Per-pattern renderers (p01-p08). Each reads SPEC[pid]['layout'] and draws to an Image. ---

def render_p01(out_path: Path) -> None:
    """P1: UI/画面レビュー型 — 勤怠アプリ画面 SS + 赤矢印 4 本 + 吹き出し 4 個."""
    spec = SPEC["p01"]
    layout = spec["layout"]
    img = _new_canvas()
    draw = ImageDraw.Draw(img)
    body_y = _draw_slide_title(draw, spec["title"], subtitle=layout["subtitle"])

    # Screenshot card (left ~60% of canvas) — taller for v2 to fit 8 rows + tabs + summary.
    card = (40, body_y + 10, 1000, body_y + 720)
    _draw_screenshot_card(draw, card, title=layout["app_title"])
    # Header inside card: user name + 3 right-side icons.
    hdr_y = body_y + 18
    # Right-side icons (drawn from right edge inwards: ⏻ 設定 通知)
    icons = layout["header_right_icons"]
    icon_x = card[2] - 14
    for icon, badge in reversed(icons):
        iw, _ = _tsize(draw, icon, _bold_font(15))
        icon_x -= (iw + 18)
        draw.text((icon_x, hdr_y - 1), icon, fill=COLORS["header_text"], font=_bold_font(15))
        if badge:
            # Red notification badge
            bx = icon_x + iw + 2
            by = hdr_y - 5
            draw.ellipse([bx, by, bx + 18, by + 18], fill=COLORS["danger"])
            draw.text((bx + 5, by + 1), badge, fill=COLORS["header_text"], font=_bold_font(11))
    # User name (left of icons)
    hdr_user_w, _ = _tsize(draw, layout["header_user"], _font(13))
    draw.text((icon_x - hdr_user_w - 18, hdr_y), layout["header_user"],
              fill=COLORS["header_text"], font=_font(13))

    # Tab bar (just below dark header)
    tab_y = body_y + 50
    tab_x = card[0] + 16
    tabs = layout["tab_bar"]["tabs"]
    active = layout["tab_bar"]["active"]
    for tab in tabs:
        tw, _ = _tsize(draw, tab, _bold_font(13))
        tab_w = tw + 24
        is_active = (tab == active)
        if is_active:
            draw.rectangle([tab_x, tab_y, tab_x + tab_w, tab_y + 30],
                           fill=COLORS["primary"], outline=COLORS["primary_dk"], width=1)
            draw.text((tab_x + 12, tab_y + 7), tab, fill=COLORS["header_text"], font=_bold_font(13))
        else:
            draw.text((tab_x + 12, tab_y + 7), tab, fill=COLORS["muted"], font=_font(13))
        tab_x += tab_w + 8
    # Filter dropdown (right side of tab bar)
    flt_text = f"{layout['filter_label']} {layout['filter_default']}"
    fw, _ = _tsize(draw, flt_text, _font(12))
    flt_x = card[2] - fw - 28
    draw.rectangle([flt_x - 6, tab_y + 2, card[2] - 16, tab_y + 28],
                   fill=COLORS["card_bg"], outline=COLORS["card_border"])
    draw.text((flt_x, tab_y + 8), flt_text, fill=COLORS["text"], font=_font(12))

    # Table (7 columns, 8 rows + summary)
    tbl_x1, tbl_y1 = card[0] + 16, body_y + 92
    col_widths = [150, 70, 70, 70, 70, 100, 280]  # 日付/出勤/退勤/実働/残業/ステータス/備考 = 810
    row_h = 32
    # Header row
    hx = tbl_x1
    for ci, header in enumerate(layout["table_header"]):
        draw.rectangle([hx, tbl_y1, hx + col_widths[ci], tbl_y1 + row_h],
                       fill=COLORS["grid"], outline=COLORS["card_border"])
        draw.text((hx + 8, tbl_y1 + 8), header, fill=COLORS["text"], font=_bold_font(12))
        hx += col_widths[ci]
    # Data rows
    for r, row in enumerate(layout["table_rows"]):
        ry = tbl_y1 + row_h * (r + 1)
        rx = tbl_x1
        bg = COLORS["bg"] if r % 2 == 0 else COLORS["card_bg"]
        for ci, val in enumerate(row):
            draw.rectangle([rx, ry, rx + col_widths[ci], ry + row_h],
                           fill=bg, outline=COLORS["card_border"])
            draw.text((rx + 8, ry + 8), val, fill=COLORS["text"], font=_font(11))
            rx += col_widths[ci]
    # Summary row (highlighted)
    sum_y = tbl_y1 + row_h * (len(layout["table_rows"]) + 1)
    rx = tbl_x1
    for ci, val in enumerate(layout["summary_row"]):
        draw.rectangle([rx, sum_y, rx + col_widths[ci], sum_y + row_h],
                       fill=COLORS["grid"], outline=COLORS["card_border"], width=1)
        draw.text((rx + 8, sum_y + 8), val, fill=COLORS["text"], font=_bold_font(11))
        rx += col_widths[ci]

    # Action buttons (below table)
    btn_y = sum_y + row_h + 20
    bx = tbl_x1
    for label in layout["buttons"]:
        w, _ = _tsize(draw, label, _bold_font(12))
        bw = w + 22
        draw.rectangle([bx, btn_y, bx + bw, btn_y + 30],
                       fill=COLORS["primary"], outline=COLORS["primary_dk"], width=2)
        draw.text((bx + 11, btn_y + 7), label, fill=COLORS["header_text"], font=_bold_font(12))
        bx += bw + 8

    # 6 vague callouts on the right side, each pointing at a target inside the card.
    # Targets are spatially anchored to specific elements (NOT spelled out in callout text).
    callout_targets = [
        (tbl_x1 + 600, tbl_y1 + row_h * 2 + row_h // 2),   # C1 → 承認待 行 (4/2 = row index 1)
        (tbl_x1 + 280, tbl_y1 + row_h * 7 + row_h // 2),   # C2 → 短実働 行 (4/9 = row index 6, 要修正)
        (tbl_x1 + 60, btn_y + 14),                          # C3 → エクスポートボタン
        (tbl_x1 + 60, tbl_y1 + 14),                         # C4 → 日付列ヘッダ
        (tbl_x1 + 540, tbl_y1 + 14),                        # C5 → ステータス列ヘッダ (列名ソート)
        (card[0] + 80, tab_y + 14),                         # C6 → 「今月」タブ
    ]
    callout_h = 80
    callout_w = 500
    for i, (label, text) in enumerate(layout["callouts"]):
        bx1 = 1040
        by1 = body_y + 20 + i * (callout_h + 16)
        _draw_callout(draw, box_rect=(bx1, by1, bx1 + callout_w, by1 + callout_h),
                      label=label, text=text, target=callout_targets[i])

    _draw_footer(draw, layout["footer"])
    img.save(out_path, "PNG")


def render_p02(out_path: Path) -> None:
    """P2 v2: Before/After 検索画面 — 詳細フィルタ + 結果行 + メニュー + フッター + 5 vague diff arrows."""
    spec = SPEC["p02"]
    layout = spec["layout"]
    img = _new_canvas()
    draw = ImageDraw.Draw(img)
    body_y = _draw_slide_title(draw, spec["title"], subtitle=layout["subtitle"])

    gap = 30
    card_w = (CANVAS_W - 3 * gap) // 2
    card_top = body_y + 10
    card_bottom = CANVAS_H - 130  # leave 130 px below for 5 diff labels + footer
    before = (gap, card_top, gap + card_w, card_bottom)
    after = (2 * gap + card_w, card_top, 2 * gap + 2 * card_w, card_bottom)

    anchors: dict[str, dict[str, tuple[int, int]]] = {}
    for card, side_key in [(before, "before"), (after, "after")]:
        s = layout[side_key]
        _draw_screenshot_card(draw, card, title=s["title"])
        cx1 = card[0] + 12

        # Header menu bar
        menu_y = card[1] + 38
        menu_text = " | ".join(s["header_menu"])
        draw.rectangle([cx1, menu_y, card[2] - 12, menu_y + 22],
                       fill=COLORS["card_bg"], outline=COLORS["card_border"])
        draw.text((cx1 + 6, menu_y + 4), menu_text, fill=COLORS["text"], font=_font(11))

        # Search bar + button
        sy = menu_y + 32
        draw.rectangle([cx1, sy, card[2] - 100, sy + 30],
                       fill=COLORS["bg"], outline=COLORS["card_border"], width=2)
        draw.text((cx1 + 8, sy + 8), s["search_placeholder"], fill=COLORS["muted"], font=_font(11))
        draw.rectangle([card[2] - 90, sy, card[2] - 12, sy + 30],
                       fill=COLORS["primary"], outline=COLORS["primary_dk"], width=2)
        draw.text((card[2] - 84, sy + 7), s["button"], fill=COLORS["header_text"], font=_bold_font(11))

        # Sort dropdown (After only)
        sort_y = sy + 38
        if s["sort_dropdown"]:
            draw.rectangle([cx1, sort_y, card[2] - 12, sort_y + 24],
                           fill=COLORS["card_bg"], outline=COLORS["card_border"])
            draw.text((cx1 + 6, sort_y + 5), s["sort_dropdown"], fill=COLORS["text"], font=_font(10))

        # Filter section
        flt_y = sort_y + (28 if s["sort_dropdown"] else 0)
        draw.text((cx1, flt_y), "Filters:", fill=COLORS["text"], font=_bold_font(11))
        flt_y += 18
        flt_y_start = flt_y
        for cat in s["filter_categories"]:
            indent = 16 if cat.startswith("  ") else 0
            draw.rectangle([cx1 + indent + 2, flt_y + 2, cx1 + indent + 14, flt_y + 14],
                           outline=COLORS["card_border"])
            draw.text((cx1 + indent + 18, flt_y), cat.strip(), fill=COLORS["text"], font=_font(10))
            flt_y += 18
        draw.text((cx1, flt_y + 4), s["filter_price_range"], fill=COLORS["text"], font=_font(10))
        flt_y += 22
        if s["filter_brands"]:
            draw.text((cx1, flt_y), s["filter_brands"], fill=COLORS["text"], font=_font(10))
            flt_y += 20
        if s["filter_stock"]:
            draw.text((cx1, flt_y), s["filter_stock"], fill=COLORS["text"], font=_font(10))
            flt_y += 20

        # Results label + 3 result rows
        res_y = flt_y + 8
        draw.text((cx1, res_y), s["result_count_label"], fill=COLORS["text"], font=_bold_font(11))
        res_y += 22
        for name, price, meta in s["result_rows"]:
            draw.rectangle([cx1, res_y, card[2] - 12, res_y + 36],
                           fill=COLORS["card_bg"], outline=COLORS["card_border"])
            draw.text((cx1 + 6, res_y + 4), name[:46], fill=COLORS["text"], font=_font(10))
            draw.text((cx1 + 6, res_y + 20), f"{price}  /  {meta}", fill=COLORS["muted"], font=_font(9))
            res_y += 40

        # Pagination
        pag_y = res_y + 8
        draw.text((cx1, pag_y), s["pagination"], fill=COLORS["text"], font=_bold_font(11))

        # Footer links
        foot_y = card[3] - 28
        draw.text((cx1, foot_y), "  |  ".join(s["footer_links"]), fill=COLORS["muted"], font=_font(9))

        anchors[side_key] = {
            "placeholder": (cx1 + 80, sy + 14),
            "filters":     (cx1 + 80, flt_y_start + 30),
            "sort":        (cx1 + 80, sort_y + 12),
            "menu":        (cx1 + 80, menu_y + 12),
            "pagination":  (cx1 + 80, pag_y + 8),
        }

    # 5 vague diff labels + dual arrows
    diff_targets = [
        (anchors["before"]["placeholder"], anchors["after"]["placeholder"]),
        (anchors["before"]["filters"],     anchors["after"]["filters"]),
        (anchors["before"]["sort"],        anchors["after"]["sort"]),
        (anchors["before"]["menu"],        anchors["after"]["menu"]),
        (anchors["before"]["pagination"],  anchors["after"]["pagination"]),
    ]
    diff_y_start = card_bottom + 8
    diff_w = 280
    diff_h = 28
    for i, (label, text) in enumerate(layout["diffs"]):
        col = i % 2
        row = i // 2
        x = 40 + col * (diff_w + 16)
        y = diff_y_start + row * (diff_h + 4)
        draw.rectangle([x, y, x + diff_w, y + diff_h], fill=COLORS["danger_bg"], outline=COLORS["danger"], width=2)
        draw.text((x + 8, y + 7), f"{label} {text}", fill=COLORS["danger_text"], font=_bold_font(11))
        for target in diff_targets[i]:
            _arrow(draw, (x + diff_w, y + diff_h // 2), target, color=COLORS["danger"], width=1, head=8)

    _draw_footer(draw, layout["footer"])
    img.save(out_path, "PNG")


def render_p03(out_path: Path) -> None:
    """P3 v2: 5 画面操作手順 — 各 step に詳細 UI 要素を縦積み + 番号付き矢印."""
    spec = SPEC["p03"]
    layout = spec["layout"]
    img = _new_canvas()
    draw = ImageDraw.Draw(img)
    body_y = _draw_slide_title(draw, spec["title"], subtitle=layout["subtitle"])

    steps = layout["steps"]
    n = len(steps)
    margin = 30
    gap = 32
    card_w = (CANVAS_W - 2 * margin - (n - 1) * gap) // n
    card_h = 600
    y1 = body_y + 50
    y2 = y1 + card_h

    for i, step in enumerate(steps):
        x1 = margin + i * (card_w + gap)
        x2 = x1 + card_w
        _draw_screenshot_card(draw, (x1, y1, x2, y2), title=step["title"])
        # Step label badge above card
        draw.rectangle([x1 + card_w // 2 - 24, y1 - 36, x1 + card_w // 2 + 24, y1 - 8],
                       fill=COLORS["primary"], outline=COLORS["primary_dk"], width=2)
        lw, _ = _tsize(draw, step["label"], _bold_font(14))
        draw.text((x1 + card_w // 2 - lw // 2, y1 - 30), step["label"],
                  fill=COLORS["header_text"], font=_bold_font(14))

        # Body lines (vertical stack inside card)
        line_y = y1 + 48
        for line in step["lines"]:
            draw.text((x1 + 10, line_y), line, fill=COLORS["text"], font=_font(11))
            line_y += 22
        # Arrow to next card
        if i < n - 1:
            arrow_y = (y1 + y2) // 2
            _arrow(draw, (x2 + 4, arrow_y), (x2 + gap - 4, arrow_y),
                   color=COLORS["text"], width=3, head=14)

    _draw_footer(draw, layout["footer"])
    img.save(out_path, "PNG")


def render_p04(out_path: Path) -> None:
    """P4 v2: 通年ダッシュボード — 棒 12 か月 + 円 8 セグ + KPI 7 + 注釈 5 (vague)."""
    spec = SPEC["p04"]
    layout = spec["layout"]
    img = _new_canvas()
    draw = ImageDraw.Draw(img)
    body_y = _draw_slide_title(draw, spec["title"], subtitle=layout["subtitle"])

    # Bar chart (top-left): 12 months
    bx1, by1 = 30, body_y + 20
    bx2, by2 = 900, by1 + 240
    _draw_screenshot_card(draw, (bx1, by1, bx2, by2), title=layout["bar_chart"]["title"])
    bars = layout["bar_chart"]["data"]
    max_v = max(v for _, v in bars)
    chart_top = by1 + 50
    chart_bot = by2 - 32
    chart_left = bx1 + 40
    chart_right = bx2 - 16
    draw.line([(chart_left, chart_bot), (chart_right, chart_bot)], fill=COLORS["text"], width=2)
    draw.line([(chart_left, chart_top), (chart_left, chart_bot)], fill=COLORS["text"], width=2)
    bar_area_w = chart_right - chart_left - 20
    bar_w = bar_area_w // len(bars) - 6
    for i, (lbl, v) in enumerate(bars):
        x = chart_left + 12 + i * (bar_w + 6)
        h = int((v / max_v) * (chart_bot - chart_top - 14))
        draw.rectangle([x, chart_bot - h, x + bar_w, chart_bot], fill=COLORS["primary"])
        draw.text((x - 1, chart_bot - h - 16), str(v), fill=COLORS["text"], font=_font(9))
        draw.text((x + 2, chart_bot + 4), lbl, fill=COLORS["text"], font=_font(9))

    # Pie chart (top-right) — 8 wedges
    cx, cy, r = 1300, by1 + 130, 90
    pie = layout["pie_chart"]
    _draw_screenshot_card(draw, (920, by1, CANVAS_W - 30, by2), title=pie["title"])
    total = sum(v for _, v in pie["data"])
    start_angle = -90
    palette = [COLORS["primary"], COLORS["success"], COLORS["warn"], COLORS["muted"],
               COLORS["danger"], COLORS["primary_dk"], COLORS["card_border"], "#7c3aed"]
    for i, (label, v) in enumerate(pie["data"]):
        sweep = (v / total) * 360
        draw.pieslice([cx - r, cy - r, cx + r, cy + r],
                      start=start_angle, end=start_angle + sweep,
                      fill=palette[i % len(palette)], outline=COLORS["bg"], width=2)
        mid = math.radians(start_angle + sweep / 2)
        lx = cx + int((r + 24) * math.cos(mid))
        ly = cy + int((r + 24) * math.sin(mid))
        draw.text((lx - 18, ly - 6), f"{label} {v}%", fill=COLORS["text"], font=_font(10))
        start_angle += sweep

    # KPI cards: 7 cards in 2 rows (4 + 3)
    kpi_y1 = by2 + 16
    kpi_h = 110
    kpis = layout["kpi_cards"]
    cards_per_row = 4
    row_gap = 12
    col_gap = 12
    kpi_w = (CANVAS_W - 2 * 30 - (cards_per_row - 1) * col_gap) // cards_per_row
    for i, (kpi_label, kpi_val, kpi_sub) in enumerate(kpis):
        col = i % cards_per_row
        row = i // cards_per_row
        kx1 = 30 + col * (kpi_w + col_gap)
        ky = kpi_y1 + row * (kpi_h + row_gap)
        _draw_screenshot_card(draw, (kx1, ky, kx1 + kpi_w, ky + kpi_h))
        draw.text((kx1 + 12, ky + 10), kpi_label, fill=COLORS["muted"], font=_bold_font(11))
        draw.text((kx1 + 12, ky + 32), kpi_val, fill=COLORS["text"], font=_bold_font(20))
        draw.text((kx1 + 12, ky + 70), kpi_sub, fill=COLORS["success"], font=_font(10))

    # 5 vague annotation rows below KPI grid
    ann_y = kpi_y1 + (kpi_h + row_gap) * 2 + 10
    for i, (label, text) in enumerate(layout["annotations"]):
        ay = ann_y + i * 30
        draw.rectangle([30, ay, CANVAS_W - 30, ay + 24],
                       fill=COLORS["danger_bg"], outline=COLORS["danger"], width=1)
        draw.text((40, ay + 5), f"{label} {text}", fill=COLORS["danger_text"], font=_bold_font(11))

    _draw_footer(draw, layout["footer"])
    img.save(out_path, "PNG")


def render_p05(out_path: Path) -> None:
    """P5 v2: 階層ドリルダウン — 上 10 モジュール / 下 8 サブモジュール + 設定 12 行 + 2 vague callouts."""
    spec = SPEC["p05"]
    layout = spec["layout"]
    img = _new_canvas()
    draw = ImageDraw.Draw(img)
    body_y = _draw_slide_title(draw, spec["title"], subtitle=layout["subtitle"])

    # Top band: 10 modules in 2 rows of 5
    modules = layout["top_level_modules"]
    highlighted = layout["highlighted_module"]
    cols = 5
    rows = 2
    mod_band_top = body_y + 16
    mod_w = (CANVAS_W - 2 * 30 - (cols - 1) * 12) // cols
    mod_h = 56
    row_gap = 8
    hi_pos = None
    for i, name in enumerate(modules):
        r = i // cols
        c = i % cols
        mx1 = 30 + c * (mod_w + 12)
        my1 = mod_band_top + r * (mod_h + row_gap)
        is_hi = (name == highlighted)
        if is_hi:
            hi_pos = (mx1 + mod_w // 2, my1 + mod_h)
        draw.rectangle(
            [mx1, my1, mx1 + mod_w, my1 + mod_h],
            fill=COLORS["card_bg"],
            outline=COLORS["danger"] if is_hi else COLORS["card_border"],
            width=3 if is_hi else 1,
        )
        nw, _ = _tsize(draw, name, _bold_font(13))
        draw.text((mx1 + (mod_w - nw) // 2, my1 + 19), name,
                  fill=COLORS["text"], font=_bold_font(13))
    mod_band_bot = mod_band_top + 2 * mod_h + row_gap

    # Drilldown arrow from highlighted module to zoom area
    zoom_y1 = mod_band_bot + 30
    zoom_y2 = zoom_y1 + 480
    if hi_pos:
        _arrow(draw, hi_pos, (CANVAS_W // 2, zoom_y1 - 4),
               color=COLORS["danger"], width=3, head=12)

    # Zoomed submodules (left half: 8 sub-modules in 4x2 grid)
    sub_card_x2 = 700
    _draw_screenshot_card(draw, (30, zoom_y1, sub_card_x2, zoom_y2),
                          title=f"拡大: {highlighted} のサブモジュール")
    sub = layout["zoom_submodules"]
    sub_cols = 2
    sub_rows = 4
    sub_w = (sub_card_x2 - 30 - 60) // sub_cols
    sub_h = 90
    for i, name in enumerate(sub):
        r = i // sub_cols
        c = i % sub_cols
        sx1 = 50 + c * (sub_w + 20)
        sy1 = zoom_y1 + 50 + r * (sub_h + 10)
        draw.rectangle([sx1, sy1, sx1 + sub_w, sy1 + sub_h],
                       fill=COLORS["bg"], outline=COLORS["card_border"], width=1)
        nw, _ = _tsize(draw, name, _bold_font(13))
        draw.text((sx1 + (sub_w - nw) // 2, sy1 + sub_h // 2 - 8), name,
                  fill=COLORS["text"], font=_bold_font(13))

    # Config table (right half) — 12 rows × 4 cols
    tbl_x1, tbl_y1 = 720, zoom_y1
    tbl_x2 = CANVAS_W - 30
    _draw_screenshot_card(draw, (tbl_x1, tbl_y1, tbl_x2, zoom_y2), title="設定パラメータ (12 項目)")
    cfg = layout["config_table"]
    col_widths = [340, 100, 110, 110]  # = 660
    hx = tbl_x1 + 12
    hy = tbl_y1 + 40
    row_h = 32
    for ci, col in enumerate(cfg["columns"]):
        draw.rectangle([hx, hy, hx + col_widths[ci], hy + row_h],
                       fill=COLORS["grid"], outline=COLORS["card_border"])
        draw.text((hx + 6, hy + 8), col, fill=COLORS["text"], font=_bold_font(11))
        hx += col_widths[ci]
    for r, row in enumerate(cfg["rows"]):
        rx = tbl_x1 + 12
        ry = hy + row_h * (r + 1)
        bg = COLORS["bg"] if r % 2 == 0 else COLORS["card_bg"]
        for ci, val in enumerate(row):
            draw.rectangle([rx, ry, rx + col_widths[ci], ry + row_h],
                           fill=bg, outline=COLORS["card_border"])
            draw.text((rx + 6, ry + 8), val, fill=COLORS["text"], font=_font(10))
            rx += col_widths[ci]

    # 2 vague callouts (below zoom area)
    callouts = layout["callouts"]
    callout_targets = [
        (CANVAS_W // 2, zoom_y1 - 30),                   # M1 → drilldown arrow
        (tbl_x1 + 200, hy + row_h * 6),                  # M2 → middle of config table
    ]
    cy_start = zoom_y2 + 14
    for i, (label, text) in enumerate(callouts):
        x = 30 + i * 600
        draw.rectangle([x, cy_start, x + 280, cy_start + 28],
                       fill=COLORS["danger_bg"], outline=COLORS["danger"], width=2)
        draw.text((x + 10, cy_start + 7), f"{label} {text}",
                  fill=COLORS["danger_text"], font=_bold_font(11))
        _arrow(draw, (x + 280, cy_start + 14), callout_targets[i],
               color=COLORS["danger"], width=1, head=8)

    _draw_footer(draw, layout["footer"])
    img.save(out_path, "PNG")


def render_p06(out_path: Path) -> None:
    """P6 v2: レビュー反映 (赤入れ) — モック 1 + 8 セクション + 25 個の短い赤コメント + 指示線."""
    spec = SPEC["p06"]
    layout = spec["layout"]
    img = _new_canvas()
    draw = ImageDraw.Draw(img)
    body_y = _draw_slide_title(draw, spec["title"], subtitle=layout["subtitle"])

    # Left: mockup placeholder with 8 stacked sections.
    mock_x1, mock_y1 = 40, body_y + 10
    mock_x2, mock_y2 = 720, body_y + 720
    _draw_screenshot_card(draw, (mock_x1, mock_y1, mock_x2, mock_y2), title="ダッシュボード モックアップ")
    sections = layout["mockup_sections"]
    sec_h = (mock_y2 - mock_y1 - 32) // len(sections)
    for i, sec in enumerate(sections):
        sy1 = mock_y1 + 32 + i * sec_h
        sy2 = sy1 + sec_h - 4
        draw.rectangle([mock_x1 + 16, sy1, mock_x2 - 16, sy2],
                       fill=COLORS["bg"], outline=COLORS["card_border"], width=1)
        draw.text((mock_x1 + 28, sy1 + 10), sec, fill=COLORS["text"], font=_bold_font(13))

    # Right: 25 short review comment bubbles, three-column compact layout.
    comments = layout["comments"]
    comment_x1 = 740
    cols = 3
    comment_w = (CANVAS_W - comment_x1 - 30 - (cols - 1) * 8) // cols
    bubble_h = 70
    row_gap = 8
    # Map every comment to its target section so each bubble has a leader arrow.
    # Bubbles wrap row-major; mod len(sections) gives a deterministic distribution.
    for i, (label, text) in enumerate(comments):
        col = i % cols
        row = i // cols
        cx1 = comment_x1 + col * (comment_w + 8)
        cy1 = body_y + 20 + row * (bubble_h + row_gap)
        cy2 = cy1 + bubble_h
        draw.rectangle([cx1, cy1, cx1 + comment_w, cy2],
                       fill=COLORS["danger_bg"], outline=COLORS["danger"], width=2)
        draw.text((cx1 + 8, cy1 + 6), label, fill=COLORS["danger_text"], font=_bold_font(12))
        wrapped = textwrap.fill(text, width=14, max_lines=2, placeholder="…")
        draw.text((cx1 + 8, cy1 + 26), wrapped, fill=COLORS["danger_text"], font=_font(11))
        # Leader arrow from leftmost-column bubble to its target section.
        if col == 0:
            target_section_idx = (i // cols) % len(sections)
            target_y = mock_y1 + 32 + target_section_idx * sec_h + sec_h // 2
            _arrow(draw, (cx1, cy1 + bubble_h // 2), (mock_x2, target_y),
                   color=COLORS["danger"], width=1, head=6)

    _draw_footer(draw, layout["footer"])
    img.save(out_path, "PNG")


def render_p07(out_path: Path) -> None:
    """P7 v2: 混合ダッシュボードページ — 9 列×8 行 表 + 8 週棒グラフ + SS + 16 行コード + 8 箇条書き."""
    spec = SPEC["p07"]
    layout = spec["layout"]
    img = _new_canvas()
    draw = ImageDraw.Draw(img)
    body_y = _draw_slide_title(draw, spec["title"], subtitle=layout["subtitle"])

    # Top-left: 9-col 8-row table
    tbl = layout["table"]
    tx1, ty1 = 40, body_y + 10
    tx2, ty2 = 1080, body_y + 380
    _draw_screenshot_card(draw, (tx1, ty1, tx2, ty2), title=tbl["title"])
    col_widths = [70, 90, 80, 80, 80, 80, 70, 90, 90]
    hx = tx1 + 8
    hy = ty1 + 44
    row_h = 32
    for ci, col in enumerate(tbl["columns"]):
        draw.rectangle([hx, hy, hx + col_widths[ci], hy + row_h],
                       fill=COLORS["grid"], outline=COLORS["card_border"])
        draw.text((hx + 6, hy + 8), col, fill=COLORS["text"], font=_bold_font(11))
        hx += col_widths[ci]
    for r, row in enumerate(tbl["rows"]):
        rx = tx1 + 8
        ry = hy + row_h * (r + 1)
        for ci, val in enumerate(row):
            draw.rectangle([rx, ry, rx + col_widths[ci], ry + row_h],
                           fill=COLORS["bg"], outline=COLORS["card_border"])
            draw.text((rx + 6, ry + 8), val, fill=COLORS["text"], font=_font(11))
            rx += col_widths[ci]

    # Top-right: bar chart (8 weeks)
    bc = layout["bar_chart"]
    bx1, by1 = 1100, body_y + 10
    bx2, by2 = CANVAS_W - 40, body_y + 380
    _draw_screenshot_card(draw, (bx1, by1, bx2, by2), title=bc["title"])
    bars = bc["data"]
    max_v = max(v for _, v in bars)
    chart_top = by1 + 60
    chart_bot = by2 - 40
    chart_left = bx1 + 30
    chart_right = bx2 - 15
    draw.line([(chart_left, chart_bot), (chart_right, chart_bot)], fill=COLORS["text"], width=2)
    bar_area_w = chart_right - chart_left - 20
    slot_w = bar_area_w // len(bars)
    bar_w = max(int(slot_w * 0.6), 12)
    for i, (lbl, v) in enumerate(bars):
        x = chart_left + 10 + i * slot_w + (slot_w - bar_w) // 2
        h = int((v / max_v) * (chart_bot - chart_top - 24))
        draw.rectangle([x, chart_bot - h, x + bar_w, chart_bot], fill=COLORS["primary"])
        draw.text((x - 4, chart_bot - h - 16), str(v), fill=COLORS["text"], font=_font(10))
        draw.text((x + 4, chart_bot + 4), lbl, fill=COLORS["text"], font=_font(11))

    # Bottom-left: screenshot caption placeholder
    ssx1, ssy1 = 40, by2 + 16
    ssx2, ssy2 = 380, ssy1 + 380
    _draw_screenshot_card(draw, (ssx1, ssy1, ssx2, ssy2), title="スクリーンショット")
    draw.rectangle([ssx1 + 16, ssy1 + 50, ssx2 - 16, ssy2 - 60],
                   fill=COLORS["card_border"], outline=COLORS["card_border"])
    cap = textwrap.fill(layout["screenshot_caption"], width=22, max_lines=3, placeholder="…")
    draw.multiline_text((ssx1 + 16, ssy2 - 52), cap,
                        fill=COLORS["text"], font=_bold_font(12), spacing=2)

    # Bottom-center: 16-line code snippet
    cx1, cy1 = 400, by2 + 16
    cx2, cy2 = 920, ssy2
    _draw_screenshot_card(draw, (cx1, cy1, cx2, cy2), title=layout["code_snippet"]["filename"])
    draw.rectangle([cx1 + 8, cy1 + 44, cx2 - 8, cy2 - 8],
                   fill="#1e1e1e", outline=COLORS["card_border"])
    code_y = cy1 + 50
    for line in layout["code_snippet"]["lines"]:
        draw.text((cx1 + 14, code_y), line, fill="#d4d4d4", font=_mono_font(11))
        code_y += 19

    # Bottom-right: 8 bullets
    bux1, buy1 = 940, by2 + 16
    bux2, buy2 = CANVAS_W - 40, ssy2
    _draw_screenshot_card(draw, (bux1, buy1, bux2, buy2), title="主要メトリクス")
    by = buy1 + 50
    for b in layout["bullets"]:
        wrapped = textwrap.fill(f"• {b}", width=24, max_lines=2, placeholder="…")
        draw.multiline_text((bux1 + 14, by), wrapped, fill=COLORS["text"], font=_font(12), spacing=2)
        by += 38

    _draw_footer(draw, layout["footer"])
    img.save(out_path, "PNG")


def _draw_dashed_line(draw, p1, p2, color, width=2, dash=8):
    """Draw a dashed straight line by stamping segments along the path."""
    x1, y1 = p1
    x2, y2 = p2
    dx = x2 - x1
    dy = y2 - y1
    length = max(1, int((dx * dx + dy * dy) ** 0.5))
    n = length // (dash * 2)
    if n == 0:
        draw.line([p1, p2], fill=color, width=width)
        return
    for i in range(n + 1):
        t1 = (i * 2) / (2 * (n + 1))
        t2 = (i * 2 + 1) / (2 * (n + 1))
        sx = x1 + dx * t1
        sy = y1 + dy * t1
        ex = x1 + dx * t2
        ey = y1 + dy * t2
        draw.line([(sx, sy), (ex, ey)], fill=color, width=width)


def render_p08(out_path: Path) -> None:
    """P8 v2: 組織図 — 4 階層 20 ノード + 部門色 + 兼務点線."""
    spec = SPEC["p08"]
    layout = spec["layout"]
    img = _new_canvas()
    draw = ImageDraw.Draw(img)
    body_y = _draw_slide_title(draw, spec["title"], subtitle=layout["subtitle"])

    nodes = layout["nodes"]  # list of (id, level, name, role, parent, joined, headcount, dept_key)
    depts = layout["departments"]
    by_level: dict[int, list[tuple]] = {}
    for n in nodes:
        by_level.setdefault(n[1], []).append(n)

    # Level y-positions for 4 layers
    level_ys = {1: body_y + 10, 2: body_y + 170, 3: body_y + 360, 4: body_y + 580}
    node_w, node_h = 150, 130
    node_positions: dict[str, tuple[int, int, int]] = {}

    for lvl, items in sorted(by_level.items()):
        total_w = len(items) * node_w + (len(items) - 1) * 18
        start_x = max(20, (CANVAS_W - total_w) // 2)
        for i, (nid, _lvl, name, role, _parent, joined, headcount, dept_key) in enumerate(items):
            nx1 = start_x + i * (node_w + 18)
            ny1 = level_ys[lvl]
            nx2 = nx1 + node_w
            ny2 = ny1 + node_h
            _, dept_color = depts[dept_key]
            # Card body
            _draw_screenshot_card(draw, (nx1, ny1, nx2, ny2))
            # Department color band on top
            draw.rectangle([nx1 + 1, ny1 + 1, nx2 - 1, ny1 + 8], fill=dept_color, outline=dept_color)
            # Avatar placeholder
            avatar_cx = nx1 + node_w // 2
            avatar_cy = ny1 + 30
            draw.ellipse([avatar_cx - 16, avatar_cy - 16, avatar_cx + 16, avatar_cy + 16],
                         fill=COLORS["card_border"], outline=COLORS["card_border"])
            # Name
            nw, _ = _tsize(draw, name, _bold_font(12))
            draw.text((nx1 + (node_w - nw) // 2, ny1 + 50), name,
                      fill=COLORS["text"], font=_bold_font(12))
            # Role
            role_disp = role if len(role) <= 20 else role[:18] + "…"
            rw, _ = _tsize(draw, role_disp, _font(10))
            draw.text((nx1 + (node_w - rw) // 2, ny1 + 70), role_disp,
                      fill=COLORS["muted"], font=_font(10))
            # Joined / headcount
            meta = f"{joined} 入社 / 配下 {headcount}"
            mw, _ = _tsize(draw, meta, _font(10))
            draw.text((nx1 + (node_w - mw) // 2, ny1 + 90), meta,
                      fill=COLORS["muted"], font=_font(10))
            # Department label
            dept_label = depts[dept_key][0]
            dw, _ = _tsize(draw, dept_label, _bold_font(10))
            draw.text((nx1 + (node_w - dw) // 2, ny1 + 108), dept_label,
                      fill=dept_color, font=_bold_font(10))
            node_positions[nid] = (avatar_cx, ny1, ny2)

    # Draw parent→child solid lines
    for nid, _lvl, _name, _role, parent, *_rest in nodes:
        if parent and parent in node_positions:
            pcx, _py1, py2 = node_positions[parent]
            ccx, cy1, _cy2 = node_positions[nid]
            mid_y = (py2 + cy1) // 2
            draw.line([(pcx, py2), (pcx, mid_y)], fill=COLORS["text"], width=2)
            draw.line([(pcx, mid_y), (ccx, mid_y)], fill=COLORS["text"], width=2)
            draw.line([(ccx, mid_y), (ccx, cy1)], fill=COLORS["text"], width=2)

    # Dotted matrix reporting lines
    for from_id, to_id, _note in layout.get("dotted_links", []):
        if from_id in node_positions and to_id in node_positions:
            fcx, fy_top, fy_bot = node_positions[from_id]
            tcx, ty_top, ty_bot = node_positions[to_id]
            _draw_dashed_line(draw, (fcx, fy_top), (tcx, ty_bot),
                              color=COLORS["danger"], width=2, dash=6)

    _draw_footer(draw, layout["footer"])
    img.save(out_path, "PNG")


def render_png(pid: str, out_path: Path) -> None:
    """Dispatch table mapping pattern id → per-pattern renderer."""
    renderers = {
        "p01": render_p01, "p02": render_p02, "p03": render_p03, "p04": render_p04,
        "p05": render_p05, "p06": render_p06, "p07": render_p07, "p08": render_p08,
    }
    if pid not in renderers:
        raise NotImplementedError(f"render_png not implemented for {pid}")
    renderers[pid](out_path)


# -----------------------------------------------------------------------------
# PPTX rendering (python-pptx)
# -----------------------------------------------------------------------------
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt


def _px(n: float) -> Emu:
    """Pixel → EMU mapping at 96 DPI so pixel layouts port over 1:1 from PIL."""
    return Emu(int(n * 9525))


def _hex_rgb(s: str) -> RGBColor:
    return RGBColor.from_string(s.lstrip("#"))


def _pptx_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _pptx_add_box(slide, x, y, w, h, *, fill=None, outline=None, outline_w=1.0,
                  shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, _px(x), _px(y), _px(w), _px(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex_rgb(fill)
    if outline is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _hex_rgb(outline)
        shp.line.width = Pt(outline_w)
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return shp


def _pptx_add_text(slide, x, y, w, h, text, *, size_pt=11, bold=False,
                   color="#111827", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(_px(x), _px(y), _px(w), _px(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = _hex_rgb(color)
    return tb


def _pptx_add_line(slide, x1, y1, x2, y2, *, color="#111827", width_pt=2.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    _px(x1), _px(y1), _px(x2), _px(y2))
    ln.line.color.rgb = _hex_rgb(color)
    ln.line.width = Pt(width_pt)
    return ln


def _pptx_title(slide, title: str, subtitle: str | None = None):
    _pptx_add_text(slide, 32, 22, 1500, 40, title, size_pt=20, bold=True)
    if subtitle:
        _pptx_add_text(slide, 32, 60, CANVAS_W - 64, 24, subtitle,
                       size_pt=11, color=COLORS["muted"])
    _pptx_add_line(slide, 32, 86, CANVAS_W - 32, 86, color=COLORS["grid"], width_pt=1.5)


def _pptx_footer(slide, text: str) -> None:
    """Muted metadata line at the very bottom of the slide."""
    _pptx_add_text(slide, 32, CANVAS_H - 24, CANVAS_W - 64, 20, text,
                   size_pt=9, color=COLORS["muted"])


def _build_slide_p01(prs) -> None:
    slide = _pptx_blank_slide(prs)
    spec = SPEC["p01"]
    layout = spec["layout"]
    _pptx_title(slide, spec["title"], subtitle=layout["subtitle"])

    # Screenshot card (taller for v2)
    card = (40, 110, 1000, 820)
    _pptx_add_box(slide, card[0], card[1], card[2] - card[0], card[3] - card[1],
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_box(slide, card[0], card[1], card[2] - card[0], 32,
                  fill=COLORS["header_bg"], outline=COLORS["header_bg"])
    _pptx_add_text(slide, card[0] + 12, card[1] + 6, 300, 24,
                   layout["app_title"], size_pt=12, bold=True, color=COLORS["header_text"])

    # Right-side icons + user (text-only, no image; pptx text is what Copilot reads)
    icons_text = " ".join([f"{ic}{f'({bd})' if bd else ''}" for ic, bd in layout["header_right_icons"]])
    _pptx_add_text(slide, card[2] - 360, card[1] + 8, 340, 24,
                   f"{layout['header_user']}    {icons_text}",
                   size_pt=11, color=COLORS["header_text"], align=PP_ALIGN.RIGHT)

    # Tab bar
    tab_y = card[1] + 40
    tab_x = card[0] + 16
    tabs = layout["tab_bar"]["tabs"]
    active = layout["tab_bar"]["active"]
    for tab in tabs:
        tab_w = len(tab) * 14 + 24
        if tab == active:
            _pptx_add_box(slide, tab_x, tab_y, tab_w, 30,
                          fill=COLORS["primary"], outline=COLORS["primary_dk"], outline_w=1.0)
            _pptx_add_text(slide, tab_x, tab_y + 7, tab_w, 22,
                           tab, size_pt=11, bold=True, color=COLORS["header_text"], align=PP_ALIGN.CENTER)
        else:
            _pptx_add_text(slide, tab_x, tab_y + 7, tab_w, 22,
                           tab, size_pt=11, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        tab_x += tab_w + 8

    # Filter dropdown
    flt_text = f"{layout['filter_label']} {layout['filter_default']}"
    _pptx_add_box(slide, card[2] - 180, tab_y + 2, 158, 26,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=0.75)
    _pptx_add_text(slide, card[2] - 174, tab_y + 7, 148, 20,
                   flt_text, size_pt=10, color=COLORS["text"])

    # Table (7 cols, 8 rows + summary)
    tbl_x1, tbl_y1 = card[0] + 16, card[1] + 82
    col_widths = [150, 70, 70, 70, 70, 100, 280]
    row_h = 32
    hx = tbl_x1
    for ci, header in enumerate(layout["table_header"]):
        _pptx_add_box(slide, hx, tbl_y1, col_widths[ci], row_h,
                      fill=COLORS["grid"], outline=COLORS["card_border"], outline_w=0.5)
        _pptx_add_text(slide, hx + 6, tbl_y1 + 8, col_widths[ci] - 12, 20,
                       header, size_pt=10, bold=True)
        hx += col_widths[ci]
    for r, row in enumerate(layout["table_rows"]):
        ry = tbl_y1 + row_h * (r + 1)
        rx = tbl_x1
        bg = COLORS["bg"] if r % 2 == 0 else COLORS["card_bg"]
        for ci, val in enumerate(row):
            _pptx_add_box(slide, rx, ry, col_widths[ci], row_h,
                          fill=bg, outline=COLORS["card_border"], outline_w=0.5)
            _pptx_add_text(slide, rx + 6, ry + 8, col_widths[ci] - 12, 20,
                           val, size_pt=9)
            rx += col_widths[ci]
    sum_y = tbl_y1 + row_h * (len(layout["table_rows"]) + 1)
    rx = tbl_x1
    for ci, val in enumerate(layout["summary_row"]):
        _pptx_add_box(slide, rx, sum_y, col_widths[ci], row_h,
                      fill=COLORS["grid"], outline=COLORS["card_border"], outline_w=0.5)
        _pptx_add_text(slide, rx + 6, sum_y + 8, col_widths[ci] - 12, 20,
                       val, size_pt=9, bold=True)
        rx += col_widths[ci]

    # Buttons
    btn_y = sum_y + row_h + 20
    bx = tbl_x1
    for label in layout["buttons"]:
        bw = len(label) * 14 + 22
        _pptx_add_box(slide, bx, btn_y, bw, 30,
                      fill=COLORS["primary"], outline=COLORS["primary_dk"], outline_w=1.5)
        _pptx_add_text(slide, bx, btn_y + 7, bw, 22,
                       label, size_pt=11, bold=True, color=COLORS["header_text"], align=PP_ALIGN.CENTER)
        bx += bw + 8

    # 6 vague callouts on right
    callout_w = 500
    callout_h = 80
    for i, (label, text) in enumerate(layout["callouts"]):
        bx1 = 1040
        by1 = 130 + i * (callout_h + 16)
        _pptx_add_box(slide, bx1, by1, callout_w, callout_h,
                      fill=COLORS["danger_bg"], outline=COLORS["danger"], outline_w=1.5,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _pptx_add_text(slide, bx1 + 12, by1 + 8, callout_w - 24, 24,
                       label, size_pt=12, bold=True, color=COLORS["danger_text"])
        _pptx_add_text(slide, bx1 + 12, by1 + 36, callout_w - 24, 36,
                       text, size_pt=12, color=COLORS["danger_text"])
    _pptx_footer(slide, layout['footer'])


def _build_slide_p02(prs) -> None:
    slide = _pptx_blank_slide(prs)
    spec = SPEC["p02"]
    layout = spec["layout"]
    _pptx_title(slide, spec["title"], subtitle=layout["subtitle"])

    gap = 30
    card_w = (CANVAS_W - 3 * gap) // 2
    card_top = 110
    card_bottom = CANVAS_H - 130
    for offset, side_key in [(gap, "before"), (2 * gap + card_w, "after")]:
        card = (offset, card_top, offset + card_w, card_bottom)
        s = layout[side_key]
        _pptx_add_box(slide, card[0], card[1], card_w, card[3] - card[1],
                      fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
        _pptx_add_box(slide, card[0], card[1], card_w, 32,
                      fill=COLORS["header_bg"], outline=COLORS["header_bg"])
        _pptx_add_text(slide, card[0] + 12, card[1] + 6, card_w - 24, 24,
                       s["title"], size_pt=12, bold=True, color=COLORS["header_text"])
        cx1 = card[0] + 12

        # Header menu
        menu_y = card[1] + 38
        _pptx_add_box(slide, cx1, menu_y, card_w - 24, 22,
                      fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=0.5)
        _pptx_add_text(slide, cx1 + 6, menu_y + 4, card_w - 36, 20,
                       " | ".join(s["header_menu"]), size_pt=10)

        # Search bar + button
        sy = menu_y + 32
        _pptx_add_box(slide, cx1, sy, card_w - 100, 30,
                      fill=COLORS["bg"], outline=COLORS["card_border"], outline_w=1.0)
        _pptx_add_text(slide, cx1 + 6, sy + 7, card_w - 110, 22,
                       s["search_placeholder"], size_pt=10, color=COLORS["muted"])
        _pptx_add_box(slide, card[2] - 90, sy, 78, 30,
                      fill=COLORS["primary"], outline=COLORS["primary_dk"], outline_w=1.0)
        _pptx_add_text(slide, card[2] - 90, sy + 7, 78, 22,
                       s["button"], size_pt=10, bold=True, color=COLORS["header_text"], align=PP_ALIGN.CENTER)

        # Sort dropdown (After only)
        sort_y = sy + 38
        if s["sort_dropdown"]:
            _pptx_add_box(slide, cx1, sort_y, card_w - 24, 24,
                          fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=0.5)
            _pptx_add_text(slide, cx1 + 6, sort_y + 5, card_w - 36, 20,
                           s["sort_dropdown"], size_pt=9)

        flt_y = sort_y + (28 if s["sort_dropdown"] else 0)
        _pptx_add_text(slide, cx1, flt_y, 200, 18, "Filters:", size_pt=11, bold=True)
        flt_y += 18
        for cat in s["filter_categories"]:
            _pptx_add_text(slide, cx1, flt_y, card_w - 24, 16,
                           f"☐ {cat}", size_pt=9)
            flt_y += 18
        _pptx_add_text(slide, cx1, flt_y + 4, card_w - 24, 18,
                       s["filter_price_range"], size_pt=9)
        flt_y += 22
        if s["filter_brands"]:
            _pptx_add_text(slide, cx1, flt_y, card_w - 24, 18, s["filter_brands"], size_pt=9)
            flt_y += 20
        if s["filter_stock"]:
            _pptx_add_text(slide, cx1, flt_y, card_w - 24, 18, s["filter_stock"], size_pt=9)
            flt_y += 20

        # Results
        res_y = flt_y + 8
        _pptx_add_text(slide, cx1, res_y, card_w - 24, 20,
                       s["result_count_label"], size_pt=10, bold=True)
        res_y += 22
        for name, price, meta in s["result_rows"]:
            _pptx_add_box(slide, cx1, res_y, card_w - 24, 36,
                          fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=0.5)
            _pptx_add_text(slide, cx1 + 6, res_y + 4, card_w - 36, 16,
                           name, size_pt=9)
            _pptx_add_text(slide, cx1 + 6, res_y + 20, card_w - 36, 14,
                           f"{price}  /  {meta}", size_pt=8, color=COLORS["muted"])
            res_y += 40

        # Pagination
        pag_y = res_y + 8
        _pptx_add_text(slide, cx1, pag_y, card_w - 24, 22,
                       s["pagination"], size_pt=11, bold=True)

        # Footer links
        _pptx_add_text(slide, cx1, card[3] - 28, card_w - 24, 20,
                       "  |  ".join(s["footer_links"]), size_pt=8, color=COLORS["muted"])

    # 5 vague diff labels at bottom (2 columns)
    diff_y_start = card_bottom + 8
    diff_w = 280
    diff_h = 28
    for i, (label, text) in enumerate(layout["diffs"]):
        col = i % 2
        row = i // 2
        x = 40 + col * (diff_w + 16)
        y = diff_y_start + row * (diff_h + 4)
        _pptx_add_box(slide, x, y, diff_w, diff_h,
                      fill=COLORS["danger_bg"], outline=COLORS["danger"], outline_w=1.5)
        _pptx_add_text(slide, x + 8, y + 7, diff_w - 16, 18,
                       f"{label} {text}", size_pt=11, bold=True, color=COLORS["danger_text"])
    _pptx_footer(slide, layout['footer'])


def _build_slide_p03(prs) -> None:
    slide = _pptx_blank_slide(prs)
    spec = SPEC["p03"]
    layout = spec["layout"]
    _pptx_title(slide, spec["title"], subtitle=layout["subtitle"])

    steps = layout["steps"]
    n = len(steps)
    margin = 30
    gap = 32
    card_w = (CANVAS_W - 2 * margin - (n - 1) * gap) // n
    card_h = 600
    y1 = 150
    for i, step in enumerate(steps):
        x1 = margin + i * (card_w + gap)
        # Step label badge above card
        _pptx_add_box(slide, x1 + card_w // 2 - 24, y1 - 36, 48, 28,
                      fill=COLORS["primary"], outline=COLORS["primary_dk"], outline_w=1.5)
        _pptx_add_text(slide, x1 + card_w // 2 - 24, y1 - 32, 48, 22,
                       step["label"], size_pt=12, bold=True, color=COLORS["header_text"], align=PP_ALIGN.CENTER)
        # Card with header
        _pptx_add_box(slide, x1, y1, card_w, card_h,
                      fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
        _pptx_add_box(slide, x1, y1, card_w, 32,
                      fill=COLORS["header_bg"], outline=COLORS["header_bg"])
        _pptx_add_text(slide, x1 + 10, y1 + 6, card_w - 20, 22,
                       step["title"], size_pt=12, bold=True, color=COLORS["header_text"])
        # Lines stacked inside card
        line_y = y1 + 44
        for line in step["lines"]:
            _pptx_add_text(slide, x1 + 10, line_y, card_w - 20, 18,
                           line, size_pt=10)
            line_y += 22
        # Arrow to next card
        if i < n - 1:
            _pptx_add_line(slide, x1 + card_w + 4, y1 + card_h // 2,
                           x1 + card_w + gap - 4, y1 + card_h // 2,
                           color=COLORS["text"], width_pt=3.0)
    _pptx_footer(slide, layout['footer'])


def _build_slide_p04(prs) -> None:
    slide = _pptx_blank_slide(prs)
    spec = SPEC["p04"]
    layout = spec["layout"]
    _pptx_title(slide, spec["title"], subtitle=layout["subtitle"])

    # Bar chart (top-left): 12 months
    bx1, by1 = 30, 110
    bx2, by2 = 900, by1 + 240
    _pptx_add_box(slide, bx1, by1, bx2 - bx1, by2 - by1,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_text(slide, bx1 + 16, by1 + 8, 500, 24,
                   layout["bar_chart"]["title"], size_pt=12, bold=True)
    bars = layout["bar_chart"]["data"]
    max_v = max(v for _, v in bars)
    chart_top, chart_bot = by1 + 50, by2 - 32
    chart_left, chart_right = bx1 + 40, bx2 - 16
    bar_area_w = chart_right - chart_left - 20
    bar_w = bar_area_w // len(bars) - 6
    for i, (lbl, v) in enumerate(bars):
        x = chart_left + 12 + i * (bar_w + 6)
        h = int((v / max_v) * (chart_bot - chart_top - 14))
        _pptx_add_box(slide, x, chart_bot - h, bar_w, h,
                      fill=COLORS["primary"], outline=COLORS["primary"])
        _pptx_add_text(slide, x - 4, chart_bot - h - 18, bar_w + 8, 16,
                       str(v), size_pt=8, align=PP_ALIGN.CENTER)
        _pptx_add_text(slide, x, chart_bot + 4, bar_w, 16,
                       lbl, size_pt=8, align=PP_ALIGN.CENTER)

    # Pie chart legend (text-only — covers all 8 segments as legend text)
    _pptx_add_box(slide, 920, by1, CANVAS_W - 950, by2 - by1,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_text(slide, 936, by1 + 8, 600, 24,
                   layout["pie_chart"]["title"], size_pt=12, bold=True)
    palette = [COLORS["primary"], COLORS["success"], COLORS["warn"], COLORS["muted"],
               COLORS["danger"], COLORS["primary_dk"], COLORS["card_border"], "#7c3aed"]
    for i, (lbl, v) in enumerate(layout["pie_chart"]["data"]):
        _pptx_add_box(slide, 936, by1 + 44 + i * 24, 18, 18,
                      fill=palette[i % len(palette)], outline=palette[i % len(palette)])
        _pptx_add_text(slide, 962, by1 + 46 + i * 24, 600, 22,
                       f"{lbl}  {v}%", size_pt=11)

    # KPI cards: 7 cards in 2 rows (4 + 3)
    kpi_y1 = by2 + 16
    kpi_h = 110
    cards_per_row = 4
    row_gap = 12
    col_gap = 12
    kpi_w = (CANVAS_W - 2 * 30 - (cards_per_row - 1) * col_gap) // cards_per_row
    for i, (kpi_label, kpi_val, kpi_sub) in enumerate(layout["kpi_cards"]):
        col = i % cards_per_row
        row = i // cards_per_row
        kx1 = 30 + col * (kpi_w + col_gap)
        ky = kpi_y1 + row * (kpi_h + row_gap)
        _pptx_add_box(slide, kx1, ky, kpi_w, kpi_h,
                      fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
        _pptx_add_text(slide, kx1 + 12, ky + 10, kpi_w - 24, 18,
                       kpi_label, size_pt=11, bold=True, color=COLORS["muted"])
        _pptx_add_text(slide, kx1 + 12, ky + 32, kpi_w - 24, 30,
                       kpi_val, size_pt=18, bold=True)
        _pptx_add_text(slide, kx1 + 12, ky + 70, kpi_w - 24, 22,
                       kpi_sub, size_pt=10, color=COLORS["success"])

    # 5 annotation rows
    ann_y = kpi_y1 + (kpi_h + row_gap) * 2 + 10
    for i, (label, text) in enumerate(layout["annotations"]):
        ay = ann_y + i * 30
        _pptx_add_box(slide, 30, ay, CANVAS_W - 60, 24,
                      fill=COLORS["danger_bg"], outline=COLORS["danger"], outline_w=1.0)
        _pptx_add_text(slide, 40, ay + 5, CANVAS_W - 80, 20,
                       f"{label} {text}", size_pt=11, bold=True, color=COLORS["danger_text"])
    _pptx_footer(slide, layout['footer'])


def _build_slide_p05(prs) -> None:
    slide = _pptx_blank_slide(prs)
    spec = SPEC["p05"]
    layout = spec["layout"]
    _pptx_title(slide, spec["title"], subtitle=layout["subtitle"])

    modules = layout["top_level_modules"]
    highlighted = layout["highlighted_module"]
    cols = 5
    mod_band_top = 110
    mod_w = (CANVAS_W - 2 * 30 - (cols - 1) * 12) // cols
    mod_h = 56
    row_gap = 8
    hi_pos = None
    for i, name in enumerate(modules):
        r = i // cols
        c = i % cols
        mx1 = 30 + c * (mod_w + 12)
        my1 = mod_band_top + r * (mod_h + row_gap)
        is_hi = (name == highlighted)
        if is_hi:
            hi_pos = (mx1 + mod_w // 2, my1 + mod_h)
        _pptx_add_box(slide, mx1, my1, mod_w, mod_h,
                      fill=COLORS["card_bg"],
                      outline=COLORS["danger"] if is_hi else COLORS["card_border"],
                      outline_w=3.0 if is_hi else 1.0)
        _pptx_add_text(slide, mx1, my1 + 18, mod_w, 22,
                       name, size_pt=12, bold=True, align=PP_ALIGN.CENTER)
    mod_band_bot = mod_band_top + 2 * mod_h + row_gap

    # Drilldown arrow
    zoom_y1 = mod_band_bot + 30
    zoom_y2 = zoom_y1 + 480
    if hi_pos:
        _pptx_add_line(slide, hi_pos[0], hi_pos[1], CANVAS_W // 2, zoom_y1 - 4,
                       color=COLORS["danger"], width_pt=3.0)

    # Zoomed submodules (left, 4x2)
    sub_card_x2 = 700
    _pptx_add_box(slide, 30, zoom_y1, sub_card_x2 - 30, zoom_y2 - zoom_y1,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_box(slide, 30, zoom_y1, sub_card_x2 - 30, 32,
                  fill=COLORS["header_bg"], outline=COLORS["header_bg"])
    _pptx_add_text(slide, 42, zoom_y1 + 6, sub_card_x2 - 50, 22,
                   f"拡大: {highlighted} のサブモジュール",
                   size_pt=12, bold=True, color=COLORS["header_text"])
    sub = layout["zoom_submodules"]
    sub_cols = 2
    sub_w = (sub_card_x2 - 30 - 60) // sub_cols
    sub_h = 90
    for i, name in enumerate(sub):
        r = i // sub_cols
        c = i % sub_cols
        sx1 = 50 + c * (sub_w + 20)
        sy1 = zoom_y1 + 50 + r * (sub_h + 10)
        _pptx_add_box(slide, sx1, sy1, sub_w, sub_h,
                      fill=COLORS["bg"], outline=COLORS["card_border"], outline_w=1.0)
        _pptx_add_text(slide, sx1, sy1 + sub_h // 2 - 10, sub_w, 22,
                       name, size_pt=12, bold=True, align=PP_ALIGN.CENTER)

    # Config table (right, 12 rows × 4 cols)
    tbl_x1 = 720
    tbl_x2 = CANVAS_W - 30
    _pptx_add_box(slide, tbl_x1, zoom_y1, tbl_x2 - tbl_x1, zoom_y2 - zoom_y1,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_box(slide, tbl_x1, zoom_y1, tbl_x2 - tbl_x1, 32,
                  fill=COLORS["header_bg"], outline=COLORS["header_bg"])
    _pptx_add_text(slide, tbl_x1 + 12, zoom_y1 + 6, 400, 22,
                   "設定パラメータ (12 項目)", size_pt=12, bold=True, color=COLORS["header_text"])
    cfg = layout["config_table"]
    col_widths = [340, 100, 110, 110]
    hx = tbl_x1 + 12
    hy = zoom_y1 + 40
    row_h = 32
    for ci, col in enumerate(cfg["columns"]):
        _pptx_add_box(slide, hx, hy, col_widths[ci], row_h,
                      fill=COLORS["grid"], outline=COLORS["card_border"], outline_w=0.5)
        _pptx_add_text(slide, hx + 6, hy + 8, col_widths[ci] - 12, 18,
                       col, size_pt=10, bold=True)
        hx += col_widths[ci]
    for r, row in enumerate(cfg["rows"]):
        rx = tbl_x1 + 12
        ry = hy + row_h * (r + 1)
        bg = COLORS["bg"] if r % 2 == 0 else COLORS["card_bg"]
        for ci, val in enumerate(row):
            _pptx_add_box(slide, rx, ry, col_widths[ci], row_h,
                          fill=bg, outline=COLORS["card_border"], outline_w=0.5)
            _pptx_add_text(slide, rx + 6, ry + 8, col_widths[ci] - 12, 18,
                           val, size_pt=9)
            rx += col_widths[ci]

    # 2 vague callouts below zoom
    cy = zoom_y2 + 14
    for i, (label, text) in enumerate(layout["callouts"]):
        x = 30 + i * 600
        _pptx_add_box(slide, x, cy, 280, 28,
                      fill=COLORS["danger_bg"], outline=COLORS["danger"], outline_w=1.5)
        _pptx_add_text(slide, x + 10, cy + 7, 270, 18,
                       f"{label} {text}", size_pt=11, bold=True, color=COLORS["danger_text"])
    _pptx_footer(slide, layout['footer'])


def _build_slide_p06(prs) -> None:
    slide = _pptx_blank_slide(prs)
    spec = SPEC["p06"]
    layout = spec["layout"]
    _pptx_title(slide, spec["title"], subtitle=layout["subtitle"])

    # Left: mockup with 8 sections
    _pptx_add_box(slide, 40, 110, 680, 720,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_box(slide, 40, 110, 680, 32,
                  fill=COLORS["header_bg"], outline=COLORS["header_bg"])
    _pptx_add_text(slide, 52, 116, 640, 24,
                   "ダッシュボード モックアップ", size_pt=12, bold=True, color=COLORS["header_text"])
    sections = layout["mockup_sections"]
    sec_h = (720 - 32) // len(sections)
    for i, sec in enumerate(sections):
        sy1 = 142 + i * sec_h
        _pptx_add_box(slide, 56, sy1, 648, sec_h - 4,
                      fill=COLORS["bg"], outline=COLORS["card_border"], outline_w=0.75)
        _pptx_add_text(slide, 68, sy1 + 10, 620, 22,
                       sec, size_pt=12, bold=True)

    # Right: 25 short comments in 3 columns
    comments = layout["comments"]
    comment_x1 = 740
    cols = 3
    comment_w = (CANVAS_W - comment_x1 - 30 - (cols - 1) * 8) // cols
    bubble_h = 70
    row_gap = 8
    for i, (label, text) in enumerate(comments):
        col = i % cols
        row = i // cols
        cx1 = comment_x1 + col * (comment_w + 8)
        cy1 = 120 + row * (bubble_h + row_gap)
        _pptx_add_box(slide, cx1, cy1, comment_w, bubble_h,
                      fill=COLORS["danger_bg"], outline=COLORS["danger"], outline_w=1.5)
        _pptx_add_text(slide, cx1 + 8, cy1 + 4, comment_w - 16, 18,
                       label, size_pt=11, bold=True, color=COLORS["danger_text"])
        _pptx_add_text(slide, cx1 + 8, cy1 + 24, comment_w - 16, 42,
                       text, size_pt=10, color=COLORS["danger_text"])
    _pptx_footer(slide, layout['footer'])


def _build_slide_p07(prs) -> None:
    slide = _pptx_blank_slide(prs)
    spec = SPEC["p07"]
    layout = spec["layout"]
    _pptx_title(slide, spec["title"], subtitle=layout["subtitle"])

    # Top-left: 9-col 8-row table
    tbl = layout["table"]
    _pptx_add_box(slide, 40, 110, 1040, 380,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_box(slide, 40, 110, 1040, 32,
                  fill=COLORS["header_bg"], outline=COLORS["header_bg"])
    _pptx_add_text(slide, 52, 116, 1000, 24, tbl["title"],
                   size_pt=12, bold=True, color=COLORS["header_text"])
    col_widths = [70, 90, 80, 80, 80, 80, 70, 90, 90]
    hy = 154
    hx = 48
    row_h = 32
    for ci, col in enumerate(tbl["columns"]):
        _pptx_add_box(slide, hx, hy, col_widths[ci], row_h,
                      fill=COLORS["grid"], outline=COLORS["card_border"])
        _pptx_add_text(slide, hx + 4, hy + 8, col_widths[ci] - 8, 20,
                       col, size_pt=10, bold=True)
        hx += col_widths[ci]
    for r, row in enumerate(tbl["rows"]):
        rx = 48
        ry = hy + row_h * (r + 1)
        for ci, val in enumerate(row):
            _pptx_add_box(slide, rx, ry, col_widths[ci], row_h,
                          fill=COLORS["bg"], outline=COLORS["card_border"])
            _pptx_add_text(slide, rx + 4, ry + 8, col_widths[ci] - 8, 20,
                           val, size_pt=10)
            rx += col_widths[ci]

    # Top-right: bar chart (8 weeks)
    bc = layout["bar_chart"]
    _pptx_add_box(slide, 1100, 110, CANVAS_W - 1140, 380,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_box(slide, 1100, 110, CANVAS_W - 1140, 32,
                  fill=COLORS["header_bg"], outline=COLORS["header_bg"])
    _pptx_add_text(slide, 1112, 116, CANVAS_W - 1180, 24, bc["title"],
                   size_pt=11, bold=True, color=COLORS["header_text"])
    bars = bc["data"]
    max_v = max(v for _, v in bars)
    chart_top, chart_bot = 170, 460
    chart_area_left = 1130
    chart_area_w = CANVAS_W - 40 - chart_area_left
    slot_w = chart_area_w // len(bars)
    bar_w = max(int(slot_w * 0.6), 12)
    for i, (lbl, v) in enumerate(bars):
        x = chart_area_left + i * slot_w + (slot_w - bar_w) // 2
        h = int((v / max_v) * (chart_bot - chart_top - 24))
        _pptx_add_box(slide, x, chart_bot - h, bar_w, h,
                      fill=COLORS["primary"], outline=COLORS["primary"])
        _pptx_add_text(slide, x - 8, chart_bot - h - 18, bar_w + 16, 18,
                       str(v), size_pt=9, align=PP_ALIGN.CENTER)
        _pptx_add_text(slide, x - 8, chart_bot + 4, bar_w + 16, 18,
                       lbl, size_pt=10, align=PP_ALIGN.CENTER)

    # Bottom-left: screenshot placeholder
    _pptx_add_box(slide, 40, 510, 340, 360,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_box(slide, 40, 510, 340, 32,
                  fill=COLORS["header_bg"], outline=COLORS["header_bg"])
    _pptx_add_text(slide, 52, 516, 320, 24,
                   "スクリーンショット", size_pt=12, bold=True, color=COLORS["header_text"])
    _pptx_add_box(slide, 56, 552, 308, 240,
                  fill=COLORS["card_border"], outline=COLORS["card_border"])
    _pptx_add_text(slide, 56, 800, 308, 60,
                   layout["screenshot_caption"], size_pt=11, bold=True)

    # Bottom-middle: 16-line code snippet
    _pptx_add_box(slide, 400, 510, 520, 360,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_box(slide, 400, 510, 520, 32,
                  fill=COLORS["header_bg"], outline=COLORS["header_bg"])
    _pptx_add_text(slide, 412, 516, 500, 24,
                   layout["code_snippet"]["filename"],
                   size_pt=12, bold=True, color=COLORS["header_text"])
    _pptx_add_box(slide, 412, 548, 496, 312,
                  fill="#1e1e1e", outline=COLORS["card_border"])
    code_y = 554
    for line in layout["code_snippet"]["lines"]:
        _pptx_add_text(slide, 420, code_y, 480, 18,
                       line, size_pt=10, color="#d4d4d4")
        code_y += 19

    # Bottom-right: 8 bullets
    _pptx_add_box(slide, 940, 510, CANVAS_W - 980, 360,
                  fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
    _pptx_add_box(slide, 940, 510, CANVAS_W - 980, 32,
                  fill=COLORS["header_bg"], outline=COLORS["header_bg"])
    _pptx_add_text(slide, 952, 516, 400, 24,
                   "主要メトリクス", size_pt=12, bold=True, color=COLORS["header_text"])
    by = 552
    for b in layout["bullets"]:
        _pptx_add_text(slide, 960, by, CANVAS_W - 1020, 36,
                       f"• {b}", size_pt=11)
        by += 38
    _pptx_footer(slide, layout['footer'])


def _build_slide_p08(prs) -> None:
    slide = _pptx_blank_slide(prs)
    spec = SPEC["p08"]
    layout = spec["layout"]
    _pptx_title(slide, spec["title"], subtitle=layout["subtitle"])

    nodes = layout["nodes"]
    depts = layout["departments"]
    by_level: dict[int, list[tuple]] = {}
    for nd in nodes:
        by_level.setdefault(nd[1], []).append(nd)
    level_ys = {1: 110, 2: 280, 3: 470, 4: 690}
    node_w, node_h = 150, 130
    node_positions: dict[str, tuple[int, int, int]] = {}

    for lvl, items in sorted(by_level.items()):
        total_w = len(items) * node_w + (len(items) - 1) * 18
        start_x = max(20, (CANVAS_W - total_w) // 2)
        for i, (nid, _lvl, name, role, _parent, joined, headcount, dept_key) in enumerate(items):
            nx1 = start_x + i * (node_w + 18)
            ny1 = level_ys[lvl]
            _, dept_color = depts[dept_key]
            # Card body
            _pptx_add_box(slide, nx1, ny1, node_w, node_h,
                          fill=COLORS["card_bg"], outline=COLORS["card_border"], outline_w=1.5)
            # Department color band on top
            _pptx_add_box(slide, nx1, ny1, node_w, 8,
                          fill=dept_color, outline=dept_color)
            # Avatar circle
            _pptx_add_box(slide, nx1 + node_w // 2 - 16, ny1 + 14, 32, 32,
                          fill=COLORS["card_border"], outline=COLORS["card_border"],
                          shape=MSO_SHAPE.OVAL)
            _pptx_add_text(slide, nx1, ny1 + 48, node_w, 20,
                           name, size_pt=11, bold=True, align=PP_ALIGN.CENTER)
            _pptx_add_text(slide, nx1, ny1 + 68, node_w, 20,
                           role, size_pt=9, color=COLORS["muted"], align=PP_ALIGN.CENTER)
            _pptx_add_text(slide, nx1, ny1 + 88, node_w, 18,
                           f"{joined} 入社 / 配下 {headcount}", size_pt=9,
                           color=COLORS["muted"], align=PP_ALIGN.CENTER)
            _pptx_add_text(slide, nx1, ny1 + 106, node_w, 18,
                           depts[dept_key][0], size_pt=9, bold=True,
                           color=dept_color, align=PP_ALIGN.CENTER)
            node_positions[nid] = (nx1 + node_w // 2, ny1, ny1 + node_h)

    for nid, _lvl, _name, _role, parent, *_rest in nodes:
        if parent and parent in node_positions:
            pcx, _, py2 = node_positions[parent]
            ccx, cy1, _ = node_positions[nid]
            mid_y = (py2 + cy1) // 2
            _pptx_add_line(slide, pcx, py2, pcx, mid_y, color=COLORS["text"], width_pt=1.75)
            _pptx_add_line(slide, pcx, mid_y, ccx, mid_y, color=COLORS["text"], width_pt=1.75)
            _pptx_add_line(slide, ccx, mid_y, ccx, cy1, color=COLORS["text"], width_pt=1.75)

    # Dotted matrix reporting lines (rendered as red lines; pptx connector lacks
    # a clean cross-platform dash style, so use color + thinner width as the
    # visual cue instead. Kept consistent with PIL output's red dashed style.)
    for from_id, to_id, _note in layout.get("dotted_links", []):
        if from_id in node_positions and to_id in node_positions:
            fcx, fy_top, fy_bot = node_positions[from_id]
            tcx, ty_top, ty_bot = node_positions[to_id]
            _pptx_add_line(slide, fcx, fy_top, tcx, ty_bot,
                           color=COLORS["danger"], width_pt=1.25)

    _pptx_footer(slide, layout['footer'])


def render_pptx(out_path: Path) -> None:
    """Build a single 8-slide PPTX, one slide per pattern p01..p08."""
    prs = Presentation()
    prs.slide_width = _px(CANVAS_W)
    prs.slide_height = _px(CANVAS_H)
    builders = {
        "p01": _build_slide_p01,
        "p02": _build_slide_p02,
        "p03": _build_slide_p03,
        "p04": _build_slide_p04,
        "p05": _build_slide_p05,
        "p06": _build_slide_p06,
        "p07": _build_slide_p07,
        "p08": _build_slide_p08,
    }
    for pid in ["p01", "p02", "p03", "p04", "p05", "p06", "p07", "p08"]:
        builders[pid](prs)
    prs.save(str(out_path))


def main() -> int:
    ap = argparse.ArgumentParser(description="Generate extraction test corpus (PPTX + 8 PNG + GT YAML)")
    ap.add_argument("--out-dir", default=str(ROOT),
                    help="Where to write artifacts (default: alongside this script)")
    args = ap.parse_args()
    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    for pid, fname in PNG_FILENAMES.items():
        path = out_dir / fname
        render_png(pid, path)
        size = path.stat().st_size
        assert size > 10_000, f"{path} is suspiciously small ({size}B) — blank canvas?"
        print(f"wrote {path} ({size} bytes)")

    pptx_path = out_dir / "extraction_test.pptx"
    render_pptx(pptx_path)
    pptx_size = pptx_path.stat().st_size
    assert pptx_size > 20_000, f"{pptx_path} is suspiciously small ({pptx_size}B)"
    print(f"wrote {pptx_path} ({pptx_size} bytes)")

    gt_path = out_dir / "ground_truth.yaml"
    emit_ground_truth_yaml(gt_path)
    gt_size = gt_path.stat().st_size
    assert gt_size > 5_000, f"{gt_path} is suspiciously small ({gt_size}B)"
    print(f"wrote {gt_path} ({gt_size} bytes)")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
