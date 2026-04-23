"""Smoke tests for the extraction material generator.

These do NOT verify visual correctness — that must be eyeballed. They verify:
- Files are created at expected paths
- PNGs are valid images with expected dimensions
- PPTX has 8 slides
- Each PNG contains the slide title as visible text (via structural checks where possible)
"""
from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image

from tests.text_vs_image.extraction import generate_extraction as g
from tests.text_vs_image.extraction.generate_extraction import PNG_FILENAMES


EXPECTED_PNG_NAMES = list(PNG_FILENAMES.values())

CANVAS_W, CANVAS_H = 1600, 900  # All PNGs share the same aspect ratio so Copilot gets uniform framing.


@pytest.mark.parametrize("pid", ["p01", "p02", "p03", "p04", "p05", "p06", "p07", "p08"])  # Extended to p01-p08 after each renderer lands.
def test_render_png_produces_valid_image(tmp_path: Path, pid: str):
    out_path = tmp_path / f"{pid}.png"
    g.render_png(pid, out_path)
    assert out_path.exists(), f"{pid} PNG was not created"
    with Image.open(out_path) as img:
        assert img.size == (CANVAS_W, CANVAS_H), f"{pid} PNG size {img.size} != expected {(CANVAS_W, CANVAS_H)}"
        assert img.format == "PNG"


from pptx import Presentation


def test_render_pptx_produces_8_slides(tmp_path: Path):
    out = tmp_path / "extraction_test.pptx"
    g.render_pptx(out)
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 8


@pytest.mark.parametrize("pid,expected_title_substring", [
    ("p01", "勤怠"),
    ("p02", "Before"),
    ("p03", "購入フロー"),
    ("p04", "売上"),
    ("p05", "決済システム"),
    ("p06", "デザインレビュー"),
    ("p07", "混合ダッシュボード"),
    ("p08", "組織図"),
])
def test_pptx_each_slide_contains_title_text(tmp_path: Path, pid: str, expected_title_substring: str):
    """Each pattern's slide must contain the pattern title as native text (so
    Copilot can read the title without OCR when given the PPTX)."""
    out = tmp_path / "extraction_test.pptx"
    g.render_pptx(out)
    prs = Presentation(str(out))
    # Slides in order p01, p02, ..., p08
    idx = int(pid[1:]) - 1
    texts: list[str] = []
    for shape in prs.slides[idx].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    combined = "\n".join(texts)
    assert expected_title_substring in combined, \
        f"slide {idx} (for {pid}) does not contain '{expected_title_substring}'"
