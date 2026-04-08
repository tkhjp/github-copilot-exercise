#!/usr/bin/env python
"""Generate PoC sample files (png, pptx, docx) for the image-describer tests.

Run once:
    python samples/generate_samples.py
"""
from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.shared import Inches as DocxInches
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt

SAMPLES = Path(__file__).resolve().parent


def _font(size: int):
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except OSError:
            continue
    return ImageFont.load_default()


def make_diagram(path: Path) -> None:
    """A simple 3-node flowchart PNG with labels."""
    w, h = 800, 400
    img = Image.new("RGB", (w, h), "white")
    draw = ImageDraw.Draw(img)
    font = _font(22)
    title_font = _font(28)

    draw.text((20, 10), "Order Processing Flow", fill="black", font=title_font)

    boxes = [
        (50, 140, 230, 240, "Receive\nOrder", "#d0e7ff"),
        (310, 140, 490, 240, "Validate\nPayment", "#fff1c0"),
        (570, 140, 750, 240, "Ship\nProduct", "#c8e6c9"),
    ]
    for x1, y1, x2, y2, label, fill in boxes:
        draw.rectangle([x1, y1, x2, y2], fill=fill, outline="black", width=3)
        lines = label.split("\n")
        for i, line in enumerate(lines):
            bbox = draw.textbbox((0, 0), line, font=font)
            tw = bbox[2] - bbox[0]
            th = bbox[3] - bbox[1]
            cx = (x1 + x2) / 2 - tw / 2
            cy = (y1 + y2) / 2 - (len(lines) * th) / 2 + i * th
            draw.text((cx, cy), line, fill="black", font=font)

    # arrows between boxes
    for start_x, end_x in [(230, 310), (490, 570)]:
        y = 190
        draw.line([(start_x, y), (end_x, y)], fill="black", width=3)
        draw.polygon(
            [(end_x, y), (end_x - 12, y - 8), (end_x - 12, y + 8)],
            fill="black",
        )

    draw.text(
        (20, 330),
        "Success rate: 97.3% (Q1 2026)",
        fill="#444444",
        font=font,
    )
    img.save(path, "PNG")


def make_bar_chart(path: Path) -> None:
    """A simple labeled bar chart PNG."""
    w, h = 600, 400
    img = Image.new("RGB", (w, h), "white")
    draw = ImageDraw.Draw(img)
    font = _font(18)
    title_font = _font(24)

    draw.text((140, 15), "Monthly Revenue (million JPY)", fill="black", font=title_font)

    # axes
    draw.line([(80, 340), (560, 340)], fill="black", width=2)  # x
    draw.line([(80, 60), (80, 340)], fill="black", width=2)  # y

    bars = [
        ("Jan", 120, "#4a90e2"),
        ("Feb", 85, "#4a90e2"),
        ("Mar", 160, "#e27a4a"),
        ("Apr", 140, "#4a90e2"),
    ]
    bar_w = 80
    gap = 40
    max_val = 200
    for i, (label, val, color) in enumerate(bars):
        x = 110 + i * (bar_w + gap)
        height = int((val / max_val) * 260)
        y_top = 340 - height
        draw.rectangle([x, y_top, x + bar_w, 340], fill=color, outline="black")
        draw.text((x + 20, 345), label, fill="black", font=font)
        draw.text((x + 18, y_top - 22), str(val), fill="black", font=font)
    img.save(path, "PNG")


def make_pptx(path: Path, png1: Path, png2: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: title
    slide = prs.slides.add_slide(blank)
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5)).text_frame
    p = tx.paragraphs[0]
    p.text = "2026 Q1 Operations Review"
    p.font.size = Pt(40)
    sub = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1)).text_frame
    sub.paragraphs[0].text = "Sample presentation for image-describer PoC"

    # Slide 2: flow diagram embedded
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(8), Inches(0.6)).text_frame.paragraphs[0].text = (
        "Order Processing Pipeline"
    )
    slide.shapes.add_picture(str(png1), Inches(1), Inches(1.5), width=Inches(8))

    # Slide 3: bar chart embedded
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(8), Inches(0.6)).text_frame.paragraphs[0].text = (
        "Q1 Revenue by Month"
    )
    slide.shapes.add_picture(str(png2), Inches(2), Inches(1.5), width=Inches(6))

    prs.save(str(path))


def make_docx(path: Path, png1: Path) -> None:
    doc = Document()
    doc.add_heading("Operations Report – Sample", level=1)
    doc.add_paragraph(
        "This sample docx embeds a flow diagram image for the image-describer PoC."
    )
    doc.add_picture(str(png1), width=DocxInches(5.5))
    doc.add_paragraph(
        "The figure above shows a three-step pipeline: receive order, validate "
        "payment, and ship product."
    )
    doc.save(str(path))


def main() -> None:
    diagram = SAMPLES / "diagram.png"
    chart = SAMPLES / "chart.png"
    pptx_path = SAMPLES / "sample.pptx"
    docx_path = SAMPLES / "sample.docx"

    make_diagram(diagram)
    make_bar_chart(chart)
    make_pptx(pptx_path, diagram, chart)
    make_docx(docx_path, diagram)

    for p in (diagram, chart, pptx_path, docx_path):
        print(f"wrote {p.relative_to(SAMPLES.parent)} ({p.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
